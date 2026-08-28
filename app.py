# -*- coding: utf-8 -*-
"""助手工作台（本地 Web 界面版，支持多助手切换 / 大模型选择 / 对话框附件 / 历史话题 / 各助手独立知识库 RAG）

启动：python -m streamlit run app.py --server.headless true --browser.gatherUsageStats false
"""
import base64
import datetime
import io
import json
import os
import re
import subprocess
import sys
import time
import uuid
from pathlib import Path

import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI

import kb
from skill_router import detect_route, should_generate, should_edit, is_copy_only

BASE_DIR = Path(__file__).parent
PROMPTS_DIR = BASE_DIR / "prompts"
STYLES_DIR = BASE_DIR / "styles"
PENGUIN_IMG = BASE_DIR / "assets" / "penguin.jpg"

st.set_page_config(page_title="助手工作台", page_icon="🐧", layout="wide")

# ---------- 初始化 ----------
load_dotenv(BASE_DIR / ".env")

UPLOAD_TYPES = ["txt", "md", "pdf", "docx", "pptx", "xlsx"]
MAX_UPLOAD_CHARS = 8000  # 本次附件注入上下文的最大字符数
MAX_HISTORY_ROUNDS = int(os.getenv("MAX_HISTORY_ROUNDS", "10"))  # 历史对话最大轮数
MAX_HISTORY_CHARS = int(os.getenv("MAX_HISTORY_CHARS", "8000"))  # 历史对话最大字符数
HISTORY_FILE = BASE_DIR / "history_topics.json"  # 历史话题本地持久化
REVIEW_ENABLED = os.getenv("REVIEW_ENABLED", "true").strip().lower() == "true"
REVIEW_MODEL = os.getenv("REVIEWER_MODEL", "").strip()  # 留空则复用当前对话模型
REVIEW_TEMP = float(os.getenv("REVIEW_TEMPERATURE", "0.3"))  # 审查温度，低温度确保审查稳定


def _model_options():
    """可选大模型：DeepSeek 使用 OPENAI_* 配置；通义千问复用 EMBEDDING_* 密钥与地址。"""
    options = {}
    ds_key = os.getenv("OPENAI_API_KEY", "").strip()
    if ds_key:
        ds_model = os.getenv("MODEL_NAME", "gpt-4o-mini").strip()
        options[f"DeepSeek（{ds_model}）"] = (
            ds_key,
            os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1").strip(),
            ds_model,
        )
    qw_key = os.getenv("EMBEDDING_API_KEY", "").strip()
    if qw_key:
        qw_model = os.getenv("QWEN_MODEL", "qwen3.8-max").strip()
        options[f"通义千问（{qw_model}）"] = (
            qw_key,
            os.getenv(
                "EMBEDDING_BASE_URL",
                "https://dashscope.aliyuncs.com/compatible-mode/v1",
            ).strip(),
            qw_model,
        )
    return options


@st.cache_resource
def get_client(api_key: str, base_url: str):
    return OpenAI(api_key=api_key, base_url=base_url)


@st.cache_resource
def _get_kb(kb_id: str) -> kb.KnowledgeBase:
    """缓存 KnowledgeBase 实例，避免每次交互重复加载。"""
    return kb.KnowledgeBase(kb_id)


def extract_text(uploaded_file) -> str:
    """从上传文件提取纯文本，支持 txt/md/pdf/docx/pptx/xlsx。"""
    ext = uploaded_file.name.rsplit(".", 1)[-1].lower()
    if ext in ("txt", "md", "markdown", "csv", "log"):
        return uploaded_file.read().decode("utf-8", errors="ignore")
    if ext == "pdf":
        import pypdf

        reader = pypdf.PdfReader(uploaded_file)
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if ext == "docx":
        import docx

        doc = docx.Document(uploaded_file)
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    table_rows.append(" | ".join(cells))
            if table_rows:
                parts.append("[表格]\n" + "\n".join(table_rows))
        return "\n".join(parts)
    if ext == "pptx":
        from pptx import Presentation

        texts = []
        for slide in Presentation(uploaded_file).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(
                        "\n".join(p.text for p in shape.text_frame.paragraphs)
                    )
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            table_rows.append(" | ".join(cells))
                    if table_rows:
                        texts.append("[表格]\n" + "\n".join(table_rows))
        return "\n".join(t for t in texts if t.strip())
    if ext == "xlsx":
        import openpyxl

        wb = openpyxl.load_workbook(uploaded_file, read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"[工作表：{ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
        return "\n".join(lines)
    raise ValueError(f"不支持的文件类型：{ext}")


def _truncate_history(messages: list, max_rounds: int = MAX_HISTORY_ROUNDS,
                      max_chars: int = MAX_HISTORY_CHARS) -> list:
    """截断历史对话：保留最近 N 轮对话，同时限制总字符数。
    始终保留 system 消息（通过外部传入），仅截断 user/assistant 轮次。
    """
    if not messages:
        return messages

    non_system = [m for m in messages if m["role"] != "system"]
    system_msgs = [m for m in messages if m["role"] == "system"]

    # 按轮次倒序保留（一轮 = 一个 user + 一个 assistant）
    rounds = []
    i = len(non_system) - 1
    while i >= 0:
        if non_system[i]["role"] == "assistant" and i > 0 and non_system[i - 1]["role"] == "user":
            rounds.insert(0, [non_system[i - 1], non_system[i]])
            i -= 2
        else:
            rounds.insert(0, [non_system[i]])
            i -= 1
        if len(rounds) >= max_rounds:
            break

    # 限制总字符数
    result_msgs = []
    total_chars = 0
    for round_pair in reversed(rounds):
        round_chars = sum(len(m.get("content", "")) for m in round_pair)
        if total_chars + round_chars > max_chars and result_msgs:
            break
        result_msgs = round_pair + result_msgs
        total_chars += round_chars

    return system_msgs + result_msgs


def _load_topics() -> dict:
    """读取历史话题（按助手名分组）。"""
    if HISTORY_FILE.exists():
        try:
            return json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _save_topics(store: dict):
    HISTORY_FILE.write_text(
        json.dumps(store, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def _new_topic() -> dict:
    return {
        "id": uuid.uuid4().hex[:8],
        "title": "新话题",
        "created": time.time(),
        "messages": [],
    }


@st.cache_resource
def _penguin_b64() -> str:
    """企鹅形象图转 base64，供 CSS 背景与欢迎卡片内联使用。"""
    return base64.b64encode(PENGUIN_IMG.read_bytes()).decode("ascii")


def _load_css() -> str:
    """从 styles/app.css 加载主题样式，替换企鹅 base64 占位符。"""
    css_file = STYLES_DIR / "app.css"
    if css_file.exists():
        css = css_file.read_text(encoding="utf-8")
    else:
        css = "<style></style>"
    return css


def inject_cute_theme():
    """注入咕咕嘎嘎可爱风：背景图 + 欢迎卡片。"""
    b64 = _penguin_b64()
    css = _load_css().replace("__PENGUIN_B64__", b64)
    st.markdown(f"<style>{css}</style>", unsafe_allow_html=True)
    st.markdown(
        '<div class="penguin-welcome">'
        f'<img src="data:image/jpeg;base64,{b64}" alt="咕咕嘎嘎">'
        "<div><h1>咕咕嘎嘎来帮你啦～ 🐧</h1>"
        "<p>选择左侧助手与大模型，上传知识库文档后即可提问。可爱模式已开启，工作效率也要萌萌哒！</p></div></div>",
        unsafe_allow_html=True,
    )


@st.cache_resource
def _reviewer_prompt() -> str:
    """加载审查员系统提示词。"""
    f = PROMPTS_DIR / "reviewer.md"
    if f.exists():
        return f.read_text(encoding="utf-8").strip()
    return "你是一名严格的 AI 输出审查员，请对回复进行结构化质量审查并按指定格式输出。"


def _run_self_review(client, model_name: str, query: str,
                     kb_context: str, upload_context: str,
                     assistant_reply: str) -> str | None:
    """执行自我审查，返回审查结果文本；失败返回 None。"""
    reviewer_system = _reviewer_prompt()
    review_user = (
        f"【用户问题】\n{query}\n\n"
        f"【知识库检索资料】\n{kb_context or '（无）'}\n\n"
        f"【本次上传文件】\n{upload_context or '（无）'}\n\n"
        f"【助手回复】\n{assistant_reply}\n\n"
        "请按审查维度逐项评估并输出结果。"
    )
    try:
        resp = client.chat.completions.create(
            model=REVIEW_MODEL or model_name,
            temperature=REVIEW_TEMP,
            messages=[
                {"role": "system", "content": reviewer_system},
                {"role": "user", "content": review_user},
            ],
        )
        return resp.choices[0].message.content
    except Exception:
        return None


def _extract_greeting(prompt_text: str) -> tuple[str, str]:
    """提取提示词文件内嵌的开场白（<!--GREETING ... GREETING-->），返回 (系统提示词, 开场白)。"""
    m = re.search(r"<!--GREETING\s*\n(.*?)\n\s*GREETING-->", prompt_text, re.DOTALL)
    if m:
        return prompt_text[: m.start()].strip(), m.group(1).strip()
    return prompt_text.strip(), ""


def _ppt_outline(client, model_name: str, content: str) -> dict | None:
    """调用模型把内容提炼为 PPT 页面大纲（JSON，含 layout 类型），失败返回 None。"""
    system = (
        "你是资深演示设计顾问。将用户内容改写为适合演示的 PPT 大纲，只输出纯 JSON，"
        "不要 markdown 代码块或其他任何文字。格式：\n"
        '{"title":"演示标题","subtitle":"一句话说明主题",'
        '"slides":[{"type":"content","title":"页标题","bullets":["要点1","要点2"],"notes":"演讲备注"}]}\n'
        "封面页与目录页由程序自动生成，不要输出。type 可选值及字段：\n"
        "- content：普通内容页，用 bullets（3~5 条）\n"
        "- two：双栏对比/并列说明，用 left_title+left[]、right_title+right[]\n"
        "- caption：架构图文页，用 bullets（层级内容）+ caption（右侧一句关键结论）\n"
        "- compare：总结/价值对比页，用 left_title+left[]（价值）、right_title+right[]（下一步）\n"
        "要求：内容页 4~8 页；最后一页必须用 compare 做总结；每条不超过 30 字，"
        "提炼为观点式短句而非原文照搬；notes 为该页演讲提示，不超过 60 字。"
    )
    try:
        resp = client.chat.completions.create(
            model=model_name,
            temperature=0.4,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": content[:12000]},
            ],
        )
        raw = (resp.choices[0].message.content or "").strip()
        raw = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw)  # 剥离可能的代码块包裹
        data = json.loads(raw)
        if isinstance(data, dict) and data.get("slides"):
            return data
    except Exception:
        return None
    return None


def _fallback_outline(content: str) -> dict:
    """模型不可用/失败时的兜底：按 Markdown 标题与段落规则切分大纲。"""
    title, slides, cur_slide = "", [], None
    for line in content.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith("#"):
            t = line.lstrip("#").strip()
            if not title:
                title = t
                continue
            cur_slide = {"type": "content", "title": t, "bullets": [], "notes": ""}
            slides.append(cur_slide)
        else:
            if cur_slide is None:
                cur_slide = {"type": "content", "title": "内容要点", "bullets": [], "notes": ""}
                slides.append(cur_slide)
            cur_slide["bullets"].append(line.lstrip("-*• ").strip())
    for s in slides:
        s["bullets"] = [b for b in s["bullets"] if b][:6]
    if not slides:
        slides = [{"type": "content", "title": "内容要点", "bullets": [content[:120]], "notes": ""}]
    # 兜底末页用 compare 总结
    if slides and slides[-1].get("type") != "compare":
        slides.append({
            "type": "compare", "title": "总结与下一步",
            "left_title": "核心价值", "left": slides[-1]["bullets"][:3],
            "right_title": "下一步", "right": ["确认范围与优先级", "分阶段推进并持续运营"],
            "notes": "",
        })
    return {"title": title or "演示文稿", "subtitle": "", "slides": slides[:10]}


_TEMPLATE_CACHE: dict[str, bytes] = {}


def _get_template_bytes() -> bytes | None:
    """加载 PPT 模板文件，支持用户上传（st.session_state.ppt_template）或默认模板。"""
    key = st.session_state.get("_ppt_template_key", "default")
    if key in _TEMPLATE_CACHE:
        return _TEMPLATE_CACHE[key]
    # 优先使用用户上传的 me.pptx，其次是默认模板
    for name in ("me.pptx", "ppt_template.pptx"):
        tmpl_path = BASE_DIR / "assets" / name
        if tmpl_path.exists():
            data = tmpl_path.read_bytes()
            _TEMPLATE_CACHE[key] = data
            return data
    return None


def _build_pptx(outline: dict, source_title: str = "") -> bytes:
    """按大纲生成 PPT。优先使用模板（带占位符），失败时回退到空白画布。"""
    template_bytes = _get_template_bytes()
    if template_bytes is not None:
        return _build_pptx_from_template(outline, source_title, template_bytes)
    return _build_pptx_fallback(outline, source_title)


def _remove_slide(prs, slide_idx: int):
    """从演示文稿中删除指定索引的幻灯片。"""
    xml_slides = prs.slides._sldIdLst
    rels = prs.part.rels
    sldId = xml_slides[slide_idx]
    rId = sldId.rId
    xml_slides.remove(sldId)
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass


def _build_pptx_from_template(outline: dict, source_title: str, template_bytes: bytes) -> bytes:
    """基于模板生成 PPT：复制模板示例页（保留背景/视觉设计），按文本框名称填充。

    模板基页约定（me.pptx，6 页）：
    - 0 = 封面（Title Slide）
    - 1 = 目录（Title and Content）
    - 2 = 普通内容页（Title and Content）
    - 3 = 双栏对比/并列说明（Two Content）
    - 4 = 架构图文页（Content with Caption）
    - 5 = 总结/价值对比页（Comparison）
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(template_bytes))
    slides_data = outline.get("slides", [])
    title = str(outline.get("title") or source_title or "演示文稿").strip()
    subtitle = str(outline.get("subtitle") or "").strip()

    BASE = {"content": 2, "two": 3, "caption": 4, "compare": 5}
    n_base = len(prs.slides)

    new_els = []  # 新页的 sldId 元素，按期望顺序记录

    # 1) 封面
    cover = _dup_slide(prs, 0)
    _fill_by_name(cover, {
        "sample-cover-title": title,
        "sample-cover-subtitle": subtitle or "咕咕嘎嘎 PM Assistant 生成",
        "sample-cover-meta": f"产品经理助手  |  {datetime.datetime.now():%Y年%m月}",
    })
    new_els.append(prs.slides._sldIdLst[-1])

    # 2) 目录（自动汇总各内容页标题）
    toc_items = [str(s.get("title") or "").strip() for s in slides_data
                 if str(s.get("title") or "").strip()]
    if toc_items:
        toc = _dup_slide(prs, 1)
        _fill_by_name(toc, {
            "sample-agenda-title": "目录",
            "sample-agenda-body": [f"{i:02d}  {t}" for i, t in enumerate(toc_items, 1)],
        })
        new_els.append(prs.slides._sldIdLst[-1])

    # 3) 内容页（按 type 选基页）
    for s in slides_data:
        s_type = str(s.get("type") or "content").strip().lower()
        slide = _dup_slide(prs, BASE.get(s_type, 2))
        _fill_typed_slide(slide, s_type, s)
        new_els.append(prs.slides._sldIdLst[-1])
        notes = str(s.get("notes") or "").strip()
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    # 4) 重排：新页移到最前（cover → 目录 → 内容页），原示例页留在尾部
    sldIdLst = prs.slides._sldIdLst
    for el in new_els:
        sldIdLst.remove(el)
    for el in reversed(new_els):
        sldIdLst.insert(0, el)

    # 5) 删除原示例页（现在位于尾部 n_base 张）
    for i in range(len(prs.slides) - 1, len(prs.slides) - 1 - n_base, -1):
        _remove_slide(prs, i)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _clone_slide_visuals(src_slide, dest_slide):
    """把 src_slide 的形状与页级背景深拷贝到 dest_slide（保留视觉设计）。

    复制关系时跳过 notesSlide——notes 与 slide 必须一对一，共享同一 notes
    部件会导致 PowerPoint 校验失败无法打开文件；新页备注由 slide.notes_slide
    按需创建独立部件。
    """
    import copy
    from pptx.oxml.ns import qn

    notes_reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"

    # 清空目标页默认形状（只保留组属性节点）
    spTree = dest_slide.shapes._spTree
    for child in list(spTree):
        if child.tag.endswith('}nvGrpSpPr') or child.tag.endswith('}grpSpPr'):
            continue
        spTree.remove(child)

    # 复制关联关系（图片/超链接等），建立 rId 映射
    rel_map = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype == notes_reltype:
            continue
        try:
            if rel.is_external:
                rel_map[rId] = dest_slide.part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True)
            else:
                rel_map[rId] = dest_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            pass

    def _remap(el):
        for node in el.iter():
            for attr in (qn('r:embed'), qn('r:link'), qn('r:id')):
                v = node.get(attr)
                if v in rel_map:
                    node.set(attr, rel_map[v])

    # 复制形状元素
    src_spTree = src_slide.shapes._spTree
    for child in list(src_spTree):
        if child.tag.endswith('}nvGrpSpPr') or child.tag.endswith('}grpSpPr'):
            continue
        el = copy.deepcopy(child)
        _remap(el)
        spTree.append(el)

    # 复制页级背景（如有）
    bg = src_slide._element.find(qn('p:bg'))
    if bg is not None:
        el = copy.deepcopy(bg)
        _remap(el)
        dest_slide._element.insert(0, el)


def _dup_slide(prs, index: int):
    """深拷贝模板幻灯片（含图片等关联部件），保留全部视觉设计。"""
    template = prs.slides[index]
    new_slide = prs.slides.add_slide(template.slide_layout)
    _clone_slide_visuals(template, new_slide)
    return new_slide


def _set_tf_text(tf, text: str):
    """替换文本框全部文本，保留首 run 的模板格式（字体/字号/颜色）。"""
    p0 = tf.paragraphs[0]
    runs = p0.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.text = text
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def _set_tf_lines(tf, lines: list):
    """多行写入文本框：首行保留模板格式，其余行克隆首行段落格式。"""
    import copy
    from pptx.oxml.ns import qn

    lines = [str(x) for x in lines if str(x).strip()] or [""]
    p0 = tf.paragraphs[0]
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    _set_tf_text(tf, lines[0])
    base_p = copy.deepcopy(tf.paragraphs[0]._p)
    last = tf.paragraphs[0]._p
    for line in lines[1:]:
        new_p = copy.deepcopy(base_p)
        runs = new_p.findall(qn("a:r"))
        for r in runs[1:]:
            new_p.remove(r)
        if runs:
            runs[0].find(qn("a:t")).text = line
        else:  # 首段无 run 的兜底
            new_p.text = line
        last.addnext(new_p)
        last = new_p


def _fill_by_name(slide, mapping: dict):
    """按文本框名称填充内容，值为 str 单行写入、list 多行写入。"""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.name in mapping:
            v = mapping[sh.name]
            if isinstance(v, list):
                _set_tf_lines(sh.text_frame, v)
            else:
                _set_tf_text(sh.text_frame, str(v))


def _col_lines(s: dict, side: str, bullets: list) -> list:
    """取双栏某一栏的行：标题行 + 条目；LLM 未分栏时把 bullets 对半分。"""
    items = [str(x).strip() for x in (s.get(side) or []) if str(x).strip()]
    other = [str(x).strip() for x in (s.get("left" if side == "right" else "right") or [])
             if str(x).strip()]
    if not items and not other and bullets:
        half = (len(bullets) + 1) // 2
        items = bullets[:half] if side == "left" else bullets[half:]
    head = str(s.get(f"{side}_title") or "").strip()
    return ([head] + items[:5]) if head else items[:6]


def _fill_typed_slide(slide, s_type: str, s: dict):
    """按版式类型填充内容页。"""
    t = str(s.get("title") or "").strip()
    bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()]

    if s_type == "two":
        _fill_by_name(slide, {
            "sample-two-title": t,
            "sample-two-left": _col_lines(s, "left", bullets),
            "sample-two-right": _col_lines(s, "right", bullets),
        })
    elif s_type == "caption":
        _fill_by_name(slide, {
            "sample-architecture-title": t,
            "sample-architecture-body": bullets[:6],
            "sample-architecture-caption": str(s.get("caption") or "关键结论").strip(),
        })
    elif s_type == "compare":
        left = [str(x).strip() for x in (s.get("left") or []) if str(x).strip()] or bullets[:3]
        right = [str(x).strip() for x in (s.get("right") or []) if str(x).strip()] or bullets[3:6]
        _fill_by_name(slide, {
            "sample-summary-title": t,
            "sample-summary-left-title": str(s.get("left_title") or "核心价值").strip(),
            "sample-summary-left-body": left[:5],
            "sample-summary-right-title": str(s.get("right_title") or "下一步").strip(),
            "sample-summary-right-body": right[:5],
        })
    else:  # content
        _fill_by_name(slide, {
            "sample-content-title": t,
            "sample-content-body": bullets[:6],
        })


def _build_pptx_fallback(outline: dict, source_title: str = "") -> bytes:
    """兜底：在空白画布上硬画（模板不可用时）。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    BG = RGBColor(0xFD, 0xF6, 0xEC)
    ACCENT = RGBColor(0xF2, 0xB9, 0x5C)
    DARK = RGBColor(0x33, 0x38, 0x3F)
    GRAY = RGBColor(0x8A, 0x7A, 0x5F)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def _new_slide():
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        return s

    title = str(outline.get("title") or source_title or "演示文稿").strip()
    cover = _new_slide()
    tb = cover.shapes.add_textbox(Inches(1.2), Inches(2.7), Inches(10.9), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, DARK
    sub = cover.shapes.add_textbox(Inches(1.25), Inches(4.4), Inches(10.9), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "🐧 由咕咕嘎嘎助手生成"
    sp.font.size, sp.font.color.rgb = Pt(16), GRAY

    for i, s in enumerate(outline.get("slides", []), 1):
        slide = _new_slide()
        title_tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
        tf = title_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i}. {str(s.get('title') or '').strip()}"
        p.font.size, p.font.bold, p.font.color.rgb = Pt(28), True, DARK
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.45), Inches(1.6), Inches(0.07)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()][:6]
        body = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.9))
        tf = body.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = f"• {b}"
            para.font.size, para.font.color.rgb = Pt(20), DARK
            para.space_after = Pt(12)
        notes = str(s.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_docx(topic_title: str, messages: list) -> bytes:
    """把当前话题导出为 Word 文档。

    遵循 anthropics/skills docx 最佳实践：
    - 使用内置 Heading 样式（非自定义），确保目录自动生成
    - 列表使用 numbering/bullet 而非手动 • 字符
    - 页边距、页码、页眉页脚完整
    - display 为空时回退到 content 字段
    - 空内容兜底，确保文档不空白
    """
    from docx import Document
    from docx.shared import Inches, Pt, Emu, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # --- 页面设置：A4，标准边距（python-docx 长度单位为 EMU，用 Cm 换算）---
    section = doc.sections[0]
    section.page_width = Cm(21.0)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.54)
    section.bottom_margin = Cm(2.54)
    section.left_margin = Cm(2.54)
    section.right_margin = Cm(2.54)

    # --- 中文字体：为默认样式设置东亚字体，避免 WPS/Word 显示乱码方块 ---
    for style_name in ("Normal", "Title", "Heading 1", "Heading 2", "List Bullet"):
        try:
            _set_style_cjk(doc.styles[style_name])
        except KeyError:
            pass

    # --- 页眉 ---
    header = section.header
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.add_run(topic_title or "对话记录")
    _set_run_cjk(run, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))

    # --- 页脚 + 页码 ---
    footer = section.footer
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run1 = fp.add_run("第 ")
    _set_run_cjk(run1, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))
    fld_begin = OxmlElement('w:fldChar')
    fld_begin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText')
    instr.text = 'PAGE'
    fld_end = OxmlElement('w:fldChar')
    fld_end.set(qn('w:fldCharType'), 'end')
    run2 = fp.add_run()
    run2._r.append(fld_begin)
    run2._r.append(instr)
    run2._r.append(fld_end)
    _set_run_cjk(run2, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))
    run3 = fp.add_run(" 页 / 共 ")
    _set_run_cjk(run3, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))
    run4 = fp.add_run()
    fld_begin2 = OxmlElement('w:fldChar')
    fld_begin2.set(qn('w:fldCharType'), 'begin')
    instr2 = OxmlElement('w:instrText')
    instr2.text = 'NUMPAGES'
    fld_end2 = OxmlElement('w:fldChar')
    fld_end2.set(qn('w:fldCharType'), 'end')
    run4._r.append(fld_begin2)
    run4._r.append(instr2)
    run4._r.append(fld_end2)
    _set_run_cjk(run4, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))
    run5 = fp.add_run(" 页")
    _set_run_cjk(run5, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))

    # --- 文档标题 ---
    doc_title = topic_title or "对话记录"
    title_h = doc.add_heading(doc_title, level=0)
    for run in title_h.runs:
        _set_run_cjk(run, color=RGBColor(0x1E, 0x3A, 0x5F))

    # --- 生成日期 ---
    date_p = doc.add_paragraph()
    date_run = date_p.add_run(f"生成日期：{datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}")
    _set_run_cjk(date_run, size=Pt(10), color=RGBColor(0x8A, 0x7A, 0x5F))

    # --- 分隔线 ---
    sep = doc.add_paragraph()
    sep_run = sep.add_run("-" * 60)
    _set_run_cjk(sep_run, size=Pt(6), color=RGBColor(0xF2, 0xB9, 0x5C))

    # --- 对话内容 ---
    content_written = False
    body_font = RGBColor(0x33, 0x33, 0x33)

    for m in messages:
        if m.get("greeting"):
            continue

        # 健壮的内容提取：display → content → 空
        text = (m.get("display") or m.get("content") or "").strip()
        if not text:
            continue

        content_written = True
        role = m["role"]

        # 角色分隔
        separator = doc.add_paragraph()
        sep_run2 = separator.add_run("-" * 40)
        _set_run_cjk(sep_run2, size=Pt(6), color=RGBColor(0xD0, 0xD0, 0xD0))

        # 角色标签（不用 emoji，避免无 emoji 字体时显示方块）
        if role == "user":
            h = doc.add_heading("【用户提问】", level=2)
        else:
            h = doc.add_heading("【助手回答】", level=2)
        for run in h.runs:
            _set_run_cjk(run, color=RGBColor(0x1E, 0x3A, 0x5F))

        # 正文段落
        for line in text.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            if _is_list_item(stripped):
                p = doc.add_paragraph(stripped, style="List Bullet")
            else:
                p = doc.add_paragraph(stripped)
            for run in p.runs:
                _set_run_cjk(run, size=Pt(11), color=body_font)

    # --- 空内容兜底 ---
    if not content_written:
        doc.add_heading("文档说明", level=1)
        fallback = doc.add_paragraph(
            "本次对话暂无可导出的内容。请在对话框中发送消息后再尝试导出。"
        )
        for run in fallback.runs:
            _set_run_cjk(run, size=Pt(12), color=RGBColor(0x66, 0x66, 0x66))

    # --- 结尾 ---
    doc.add_paragraph()
    end_p = doc.add_paragraph()
    end_run = end_p.add_run("—— 由咕咕嘎嘎 PM Assistant 生成 ——")
    _set_run_cjk(end_run, size=Pt(9), color=RGBColor(0x8A, 0x7A, 0x5F))
    end_p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _set_style_cjk(style, font_name: str = "微软雅黑"):
    """为样式设置东亚字体（w:eastAsia），确保中文在 Word/WPS 正确渲染。"""
    from docx.oxml.ns import qn
    style.font.name = font_name
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.find(qn('w:rFonts'))
    if rfonts is None:
        rfonts = rpr.makeelement(qn('w:rFonts'), {})
        rpr.append(rfonts)
    rfonts.set(qn('w:eastAsia'), font_name)


def _set_run_cjk(run, font_name: str = "微软雅黑", size=None, color=None, bold=None):
    """为 run 设置中文字体（含 w:eastAsia）、字号、颜色。"""
    from docx.oxml.ns import qn
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
    if size is not None:
        run.font.size = size
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def _is_list_item(text: str) -> bool:
    """判断是否为列表项（以 -, *, •, 数字., 数字) 开头）。"""
    import re
    return bool(re.match(r'^[-*•]|\d+[.)、]', text))


def _modify_pptx(pptx_bytes: bytes, patch_ops: list[dict],
                 client=None, model_name=None, instruction: str = "") -> bytes | None:
    """两阶段 PPT 修改：hands-on-deck 结构 + python-pptx 填充新空白页。

    Phase 1: hands-on-deck 处理结构变更（增/删/复制幻灯片、替换已有文本）
    Phase 2: python-pptx 检测新增空白页，用 LLM 生成内容填充
    """
    hod_dir = BASE_DIR / "hands_on_deck"
    deck_py = hod_dir / "scripts_deck.py"
    if not deck_py.exists():
        st.session_state["last_deck_error"] = f"未找到 {deck_py}"
        return None

    import tempfile
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / "input.pptx"
        output_path = Path(tmpdir) / "output.pptx"
        patch_path = Path(tmpdir) / "patch.json"

        input_path.write_bytes(pptx_bytes)
        patch_path.write_text(json.dumps(patch_ops, ensure_ascii=False), encoding="utf-8")

        try:
            result = subprocess.run(
                [sys.executable, str(deck_py), str(input_path), "apply",
                 str(patch_path), "-o", str(output_path), "--fix"],
                capture_output=True, text=True, cwd=str(hod_dir), timeout=120
            )
        except subprocess.TimeoutExpired:
            st.session_state["last_deck_error"] = "hands-on-deck 执行超时（>120 秒）"
            return None
        if result.returncode != 0 or not output_path.exists():
            st.session_state["last_deck_error"] = (result.stderr or result.stdout or "无错误输出")[-4000:]
            return None

        modified_bytes = output_path.read_bytes()

        # Phase 2: 检测并填充新增空白页
        if client and model_name and instruction:
            modified_bytes = _post_process_new_slides(modified_bytes, client, model_name, instruction)

        return modified_bytes


def _deck_inventory(pptx_bytes: bytes) -> str:
    """调用 hands-on-deck inspect --brief 获取幻灯片形状清单，注入 patch 生成 prompt。"""
    hod_dir = BASE_DIR / "hands_on_deck"
    deck_py = hod_dir / "scripts_deck.py"
    if not deck_py.exists():
        return ""
    import tempfile
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / "input.pptx"
        input_path.write_bytes(pptx_bytes)
        try:
            result = subprocess.run(
                [sys.executable, str(deck_py), str(input_path), "inspect", "--brief"],
                capture_output=True, text=True, cwd=str(hod_dir), timeout=60
            )
        except subprocess.TimeoutExpired:
            return ""
        if result.returncode != 0:
            return ""
        return result.stdout[:8000]


def _instruction_to_patch(client, model_name, instruction: str, pptx_bytes: bytes) -> list | None:
    """将自然语言修改指令转换为 hands-on-deck patch 格式。"""
    inventory = _deck_inventory(pptx_bytes)
    inventory_block = (
        f"当前 PPT 形状清单（slide 索引从 0 开始；shape id 为 sN 编号，引用形状必须用它）：\n{inventory}\n"
        if inventory else ""
    )
    prompt = f"""你是 PPT 编辑专家。将用户的修改指令转换为 hands-on-deck patch JSON 数组。

{inventory_block}
支持的 patch 操作（op 字段，字段名严格，不得更改）：
- replace-text: 替换文本  {{"op":"replace-text","scope":"deck","from":"原文本","to":"新文本"}}
  scope 只能是 "deck"（全部幻灯片）| "slide"（单页，须同时给 "slide":N）| "master"（母版/版式）。禁止写成 "slide:N"
- set-text: 重写已有形状的文本  {{"op":"set-text","slide":0,"shape":"s5","text":["第一行","第二行"]}}
- delete: 删除形状      {{"op":"delete","slide":0,"shape":"s5"}}
- duplicate: 复制形状   {{"op":"duplicate","slide":0,"shape":"s5","offset":[0,1.2]}}
- add-slide: 添加幻灯片  {{"op":"add-slide"}}
  可选 "at":N（0-based 插入位置，省略则追加到末尾）、"layout":版式名字符串或 0-based 整数索引（省略用最空白版式）。此操作只创建空白页，内容会由后续步骤填充
- set-notes: 设置备注    {{"op":"set-notes","slide":0,"notes":"备注内容"}}

重要规则：
1. "slide" 一律是 0-based 整数；"shape" 一律是形状清单中的 id 字符串（如 "s5"），不是数字下标
2. 禁止使用 shape_idx / idx / layout_idx 字段；禁止 scope:"slide:N" 写法；set-notes 的备注字段是 "notes" 不是 "text"
3. 用户要求"添加新页/增加几页/新增总结页"时，用 add-slide（内容填充由系统后续处理）
4. 用户要求"修改/替换已有文本"时，优先用 replace-text；文本被拆分在多个 run 中导致替换失败时改用 set-text
5. 只依据上方形状清单中的信息定位形状，清单中没有的页不要猜测形状

用户指令：{instruction}

请返回纯 JSON 数组，不要添加任何解释文字。如果无法理解指令，返回空数组 []。"""

    try:
        resp = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
        )
        raw = resp.choices[0].message.content.strip()
        raw = raw.removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        patch = json.loads(raw)
        if isinstance(patch, list) and len(patch) > 0:
            return patch
        return None
    except Exception:
        return None


def _post_process_new_slides(pptx_bytes: bytes, client, model_name, instruction: str) -> bytes:
    """Phase 2: 检测 PPT 中的空白新页，用 LLM 生成内容填充。

    判定"空白页"的标准：有标题占位符但标题为空，或正文占位符无实际文本。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_count = len(prs.slides)

    # 找出空白页的索引（相对于原始文件，add-slide 追加的是最后几页）
    blank_indices = []
    for i, slide in enumerate(prs.slides):
        if _is_slide_blank(slide):
            blank_indices.append(i)

    if not blank_indices:
        return pptx_bytes

    # 样式源页：最后一张非空白页（add-slide 追加在末尾，前面都是原页）
    style_slide = None
    for i in range(slide_count - 1, -1, -1):
        if i not in blank_indices:
            style_slide = prs.slides[i]
            break

    # 用 LLM 为每个空白页生成内容；先复制原页样式再填充，避免纯空白页
    for idx in blank_indices:
        slide = prs.slides[idx]
        if style_slide is not None:
            _clone_slide_visuals(style_slide, slide)
            _clear_slide_text(slide)
        content = _llm_generate_slide_content(client, model_name, instruction, idx + 1, len(blank_indices))
        _fill_slide_content(slide, content)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _is_slide_blank(slide) -> bool:
    """判断幻灯片是否为空白（无实质文本，且无图片/表格/图表/SmartArt 等非文本内容）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    non_text_types = (
        MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE,
        MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.IGX_GRAPHIC,
    )
    for shape in slide.shapes:
        try:
            if shape.shape_type in non_text_types:
                return False
        except Exception:  # 无法识别的形状类型，保守视为非空白
            return False
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) > 1:
                return False
    return True


def _llm_generate_slide_content(client, model_name, instruction: str,
                                 slide_num: int, total_new: int) -> dict:
    """为单个空白页生成标题和内容。"""
    prompt = f"""你是 PPT 内容撰写专家。根据用户的修改指令，为第 {slide_num}/{total_new} 页生成 PPT 内容。

用户指令：{instruction}

请返回 JSON 格式：
{{"title": "页面标题", "subtitle": "副标题（可选，无则省略）", "bullets": ["要点1", "要点2", "要点3"]}}

- title: 5-15 字的简洁标题
- bullets: 3-6 个要点，每个 15-40 字
- 如果是总结/结束页，title 可用"总结"、"Thanks"等，bullets 可省略

返回纯 JSON，不要解释。"""

    try:
        resp = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
        )
        raw = resp.choices[0].message.content.strip()
        raw = raw.removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        return json.loads(raw)
    except Exception:
        return {"title": f"新增页面 {slide_num}", "bullets": ["内容待补充"]}


def _clear_slide_text(slide):
    """清空幻灯片所有文本框的文本，保留形状与字体样式。"""
    from pptx.oxml.ns import qn

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for t in shape.text_frame._txBody.iter(qn('a:t')):
            t.text = ""


def _find_title_body(slide):
    """定位标题框与正文框：优先占位符，否则按位置启发式。"""
    title_shape = None
    body_shape = None
    for shape in slide.placeholders:
        try:
            idx = shape.placeholder_format.idx
        except Exception:
            continue
        if idx == 0:
            title_shape = shape
        elif idx == 1:
            body_shape = shape
    if title_shape is not None and body_shape is not None:
        return title_shape, body_shape

    text_shapes = [sh for sh in slide.shapes if sh.has_text_frame]
    if not text_shapes:
        return None, None
    # 标题：垂直位置最高（top 最小）的文本框；正文：剩余里面积最大
    text_shapes.sort(key=lambda sh: (sh.top if sh.top is not None else 0))
    title_shape = text_shapes[0]
    rest = text_shapes[1:]
    body_shape = None
    if rest:
        body_shape = max(rest, key=lambda sh: ((sh.width or 0) * (sh.height or 0)))
    return title_shape, body_shape


def _fill_slide_content(slide, content: dict):
    """将内容填充到幻灯片，尽量复用已有文本框（保留原 PPT 的字体/字号/颜色）。

    克隆样式页后的空白页会带上原页文本框，这里复用它们而非新建，从而继承原样式。
    """
    from pptx.util import Inches

    title_text = str(content.get("title") or "").strip()
    bullets = [str(b).strip() for b in (content.get("bullets") or []) if str(b).strip()]
    subtitle = str(content.get("subtitle") or "").strip()

    title_shape, body_shape = _find_title_body(slide)

    # 标题
    if title_text:
        if title_shape is not None:
            _set_tf_text(title_shape.text_frame, title_text)
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
            txBox.text_frame.word_wrap = True
            _set_tf_text(txBox.text_frame, title_text)

    # 副标题
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.6))
        txBox.text_frame.word_wrap = True
        _set_tf_text(txBox.text_frame, subtitle)

    # 正文要点
    if bullets:
        if body_shape is not None:
            _set_tf_lines(body_shape.text_frame, bullets)
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.3))
            txBox.text_frame.word_wrap = True
            _set_tf_lines(txBox.text_frame, bullets)


def _modify_docx(docx_bytes: bytes, instruction: str, client=None, model_name=None) -> bytes | None:
    """使用 python-docx 增量编辑 DOCX（文本替换、段落删除、追加等）。"""
    try:
        from docx import Document

        doc = Document(io.BytesIO(docx_bytes))

        if client and model_name:
            ops = _llm_docx_instruction(client, model_name, instruction)
        else:
            ops = _parse_docx_instruction(instruction)

        for op in ops:
            if op["op"] == "replace":
                _docx_replace_text(doc, op["from"], op["to"])
            elif op["op"] == "delete_paragraph":
                _docx_delete_paragraphs(doc, op.get("keyword", ""))
            elif op["op"] == "append":
                style = op.get("style", "")
                p = doc.add_paragraph(op.get("text", ""))
                if style:
                    try:
                        p.style = doc.styles[style]
                    except Exception:
                        pass
            elif op["op"] == "insert_after":
                _docx_insert_after(doc, op.get("keyword", ""), op.get("text", ""))

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except Exception:
        return None


def _llm_docx_instruction(client, model_name, instruction: str) -> list:
    """用 LLM 将自然语言指令解析为 DOCX 操作列表。"""
    prompt = f"""你是 Word 文档编辑专家。将用户的修改指令转换为 JSON 操作数组。

支持的操作：
- {{"op":"replace","from":"原文本","to":"新文本"}}  — 全文替换文本
- {{"op":"delete_paragraph","keyword":"关键词"}}    — 删除包含关键词的段落
- {{"op":"append","text":"追加内容","style":"Normal"}}  — 在文档末尾追加段落
- {{"op":"insert_after","keyword":"锚点","text":"插入内容"}} — 在包含关键词的段落后插入新段落

用户指令：{instruction}

返回纯 JSON 数组，不要任何解释。"""
    try:
        resp = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
        )
        raw = resp.choices[0].message.content.strip()
        raw = raw.removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        ops = json.loads(raw)
        if isinstance(ops, list):
            return ops
    except Exception:
        pass
    return _parse_docx_instruction(instruction)


def _parse_docx_instruction(instruction: str) -> list:
    """正则兜底解析 DOCX 指令。"""
    ops = []
    import re

    for m in re.finditer(r'(?:把|将)(.+?)(?:改成|替换为|换成|改为)(.+?)(?=[，,。；;]|$)', instruction):
        ops.append({"op": "replace", "from": m.group(1).strip().strip('"“”'), "to": m.group(2).strip().strip('"“”')})

    for m in re.finditer(r'删除(?:包含|含有|关于)?(.+?)(?:的|之)?段落', instruction):
        ops.append({"op": "delete_paragraph", "keyword": m.group(1).strip()})

    for m in re.finditer(r'(?:在末尾|最后|追加|添加)(.+?)(?:内容|段落|文字)', instruction):
        ops.append({"op": "append", "text": m.group(1).strip()})

    return ops


def _docx_insert_after(doc, keyword: str, text: str):
    """在包含关键词的段落后插入新段落（继承锚点段落属性，文本只写入一次）。"""
    import copy
    from docx.text.paragraph import Paragraph

    if not keyword:
        return
    for para in doc.paragraphs:
        if keyword in para.text:
            new_p = copy.deepcopy(para._element)
            # 清掉复制出的 run/超链接等内容，只保留 pPr 段落属性
            for child in list(new_p):
                if not child.tag.endswith("}pPr"):
                    new_p.remove(child)
            para._element.addnext(new_p)
            Paragraph(new_p, para._parent).add_run(text)
            break


def _docx_replace_text(doc, old: str, new: str):
    """在 DOCX 所有段落和表格中替换文本。"""
    for para in doc.paragraphs:
        _replace_in_paragraph(para, old, new)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _replace_in_paragraph(para, old, new)


def _replace_in_paragraph(para, old: str, new: str):
    """在段落中替换文本（跨 run 处理）。"""
    full = "".join(run.text for run in para.runs)
    if old not in full:
        return
    new_text = full.replace(old, new)
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""


def _docx_delete_paragraphs(doc, keyword: str):
    """删除包含关键词的段落。"""
    if not keyword:
        return
    to_remove = []
    for i, para in enumerate(doc.paragraphs):
        if keyword in para.text:
            to_remove.append(para)
    for para in to_remove:
        p = para._element
        p.getparent().remove(p)


def _autodownload(ext: str):
    """注入 JS 自动触发一次下载点击，实现“生成完直接下载”。"""
    import streamlit.components.v1 as components

    js = (
        "const links = window.parent.document.querySelectorAll('a[download$=\"" + ext + "\"]');"
        "if (links.length) links[links.length - 1].click();"
    )
    components.html(f"<script>{js}</script>", height=0)


if "topic_store" not in st.session_state:
    st.session_state.topic_store = _load_topics()
if "cur_topic" not in st.session_state:
    st.session_state.cur_topic = {}
if "rename_topic_id" not in st.session_state:
    st.session_state.rename_topic_id = None

inject_cute_theme()

# ---------- 侧边栏 ----------
with st.sidebar:
    assistants = {}
    if PROMPTS_DIR.exists():
        assistants = {
            p.stem: p for p in sorted(PROMPTS_DIR.glob("*.md"))
            if p.stem != "reviewer"
        }
    if not assistants:
        st.error("prompts/ 目录下没有助手提示词文件（.md），请先添加。")
        st.stop()
    names = list(assistants.keys())
    if st.session_state.get("assistant_sel") not in names:
        st.session_state.pop("assistant_sel", None)
    chosen = st.selectbox("当前助手", names, key="assistant_sel")
    SYSTEM_PROMPT, GREETING = _extract_greeting(
        assistants[chosen].read_text(encoding="utf-8")
    )
    knowledge_base = _get_kb(chosen)  # 使用缓存的知识库实例

    st.title(f"🐧 {chosen}")

    # ---------- 历史话题（各助手独立，上下文互不串扰） ----------
    store = st.session_state.topic_store
    topics = store.setdefault(chosen, [])

    # 清理多余空话题（最多保留 1 个），避免刷新/进入时侧边栏堆满"新话题"
    def _is_empty(t):
        return not any(not m.get("greeting") for m in t["messages"])

    empties = [t for t in topics if _is_empty(t)]
    if len(empties) > 1:
        for t in empties[1:]:
            topics.remove(t)
        _save_topics(store)

    cur = next(
        (t for t in topics if t["id"] == st.session_state.cur_topic.get(chosen)),
        None,
    )
    if cur is None:  # 首次进入或当前话题已被删除：优先复用空话题，不重复新建
        cur = empties[0] if empties else _new_topic()
        if cur not in topics:
            topics.insert(0, cur)
        st.session_state.cur_topic[chosen] = cur["id"]
        _save_topics(store)
    messages = cur["messages"]
    if not messages and GREETING:  # 新话题由助手主动发起开场白（仅展示，不进模型上下文）
        messages.append(
            {"role": "assistant", "display": GREETING, "content": "", "greeting": True}
        )
        _save_topics(store)

    st.subheader("历史对话")
    if st.button("+ 新建话题", key="new_topic", type="tertiary", use_container_width=True):
        cur = _new_topic()
        topics.insert(0, cur)
        st.session_state.cur_topic[chosen] = cur["id"]
        _save_topics(store)
        st.rerun()

    for t in topics:
        is_renaming = st.session_state.rename_topic_id == t["id"]
        if is_renaming:
            new_title = st.text_input(
                "重命名话题",
                value=t["title"],
                key=f"rename_input_{t['id']}",
                label_visibility="collapsed",
            )
            c1, c2 = st.columns([1, 1])
            with c1:
                if st.button("✓", key=f"rename_ok_{t['id']}", type="tertiary"):
                    t["title"] = new_title[:30] if new_title.strip() else t["title"]
                    st.session_state.rename_topic_id = None
                    _save_topics(store)
                    st.rerun()
            with c2:
                if st.button("✕", key=f"rename_cancel_{t['id']}", type="tertiary"):
                    st.session_state.rename_topic_id = None
                    st.rerun()
        else:
            c1, c2, c3 = st.columns([6, 1, 1])
            with c1:
                is_active = t["id"] == cur["id"]
                label = f"{'🟢' if is_active else '💬'} {t['title']}"
                if st.button(
                    label,
                    key=f"topic_{t['id']}",
                    type="primary" if is_active else "secondary",
                    use_container_width=True,
                ):
                    st.session_state.cur_topic[chosen] = t["id"]
                    st.rerun()
            with c2:
                if st.button(
                    "✎", key=f"topic_rename_{t['id']}", type="tertiary", help="重命名"
                ):
                    st.session_state.rename_topic_id = t["id"]
                    st.rerun()
            with c3:
                if st.button(
                    "🗑", key=f"topic_del_{t['id']}", type="tertiary", help="删除"
                ):
                    topics[:] = [x for x in topics if x["id"] != t["id"]]
                    if st.session_state.cur_topic.get(chosen) == t["id"]:
                        st.session_state.cur_topic[chosen] = (
                            topics[0]["id"] if topics else None
                        )
                    if st.session_state.rename_topic_id == t["id"]:
                        st.session_state.rename_topic_id = None
                    _save_topics(store)
                    st.rerun()

    # ---------- 大模型选择（须先于编辑模式，编辑功能依赖 client） ----------
    options = _model_options()
    labels = list(options.keys())
    if st.session_state.get("model_sel") not in labels:
        st.session_state.pop("model_sel", None)
    if labels:
        label = st.selectbox("当前大模型", labels, key="model_sel")
        api_key, base_url, model_name = options[label]
        client = get_client(api_key, base_url)
        st.caption(f"当前模型：{model_name}")
    else:
        client = None
        model_name = None
        st.error("未配置任何大模型密钥：本地请编辑 .env；Streamlit Cloud 请在 Settings → Secrets 中配置后重启应用。")

    # ---------- 编辑模式（增量修改已有 PPT/DOCX，只输出副本） ----------
    st.subheader("🛠️ 编辑已有文档")
    edit_file = st.file_uploader(
        "上传要编辑的 PPT/DOCX",
        type=["pptx", "docx"],
        key="edit_target",
        help="上传后用自然语言描述修改需求，AI 生成 patch 增量编辑（永不覆盖源文件）",
    )
    edit_instruction = st.text_area(
        "修改指令",
        placeholder="例如：把「营销场景」改成「武警舆情」，在最后添加一页项目总结",
        key="edit_instr",
        height=80,
    )

    if edit_file is not None:
        file_bytes = edit_file.read()
        decision = detect_route("upload_file", edit_file.name, file_bytes)

        for w in decision.warnings:
            st.warning(w)
        for bw in decision.boundary_warnings:
            st.warning(bw)
        if not decision.boundary_warnings:
            st.caption("✅ 文件检测：未发现 SmartArt/动画/修订标记等复杂元素")

        st.caption(f"🔀 路由：**{decision.route}** | 安全：**{decision.safety['output_mode']}**")

        if decision.route == "edit_ppt" and edit_instruction and client is not None:
            if st.button("🚀 执行 PPT 修改", key="exec_edit_ppt"):
                with st.spinner("🔧 正在生成 patch 并应用到 PPT（约 10~60 秒）…"):
                    patch = _instruction_to_patch(client, model_name, edit_instruction, file_bytes)
                    if patch:
                        result = _modify_pptx(file_bytes, patch,
                                              client=client, model_name=model_name,
                                              instruction=edit_instruction)
                        if result:
                            st.session_state.edited_bytes = result
                            st.session_state.edited_name = f"edited_{edit_file.name}"
                            st.success(f"✅ 修改成功！输出 {len(result)} bytes 副本")
                        else:
                            st.error("❌ hands-on-deck 执行失败，请检查 patch 格式")
                            deck_err = st.session_state.get("last_deck_error")
                            if deck_err:
                                with st.expander("查看错误详情"):
                                    st.code(deck_err, language="text")
                    else:
                        st.error("❌ 无法生成 patch 指令，请简化修改描述")

        elif decision.route == "edit_docx":
            if edit_instruction and st.button("🚀 执行 DOCX 修改", key="exec_edit_docx"):
                with st.spinner("🔧 正在修改 DOCX…"):
                    edited = _modify_docx(file_bytes, edit_instruction, client=client, model_name=model_name)
                    if edited:
                        st.session_state.edited_bytes = edited
                        st.session_state.edited_name = f"edited_{edit_file.name}"
                        st.success(f"✅ 修改成功！输出 {len(edited)} bytes 副本")
                    else:
                        st.error("❌ DOCX 修改失败")

    if st.session_state.get("edited_bytes"):
        st.download_button(
            "⬇️ 下载编辑后的文件",
            data=st.session_state["edited_bytes"],
            file_name=st.session_state.get("edited_name", "edited_output.bin"),
            mime="application/octet-stream",
            key="edited_dl",
        )

    with st.expander("Agent 自我审查", expanded=False):
        review_enabled = st.toggle(
            "启用回复审查",
            value=REVIEW_ENABLED,
            help="每次回复后自动进行质量审查（准确性/完整性/合规性等），审查结果可展开查看",
        )
        if review_enabled:
            review_model_label = REVIEW_MODEL or "复用当前对话模型"
            st.caption(f"审查模型：{review_model_label} | 温度：{REVIEW_TEMP}")

    with st.expander("PPT 模板（可选）", expanded=False):
        tmpl = st.file_uploader(
            "上传 .pptx 模板（可选，用于生成 PPT 时套用样式）",
            type=["pptx"],
            key="ppt_template_uploader",
            help="模板需含标准占位符（标题/正文）。不上传则使用默认企鹅奶油风模板。",
        )
        if tmpl:
            st.session_state["_ppt_template_key"] = f"user_{tmpl.size}_{tmpl.name}"
            _TEMPLATE_CACHE[f"user_{tmpl.size}_{tmpl.name}"] = tmpl.read()
            st.success(f"已加载模板：{tmpl.name}")
        elif "_ppt_template_key" in st.session_state:
            if st.button("清除自定义模板", key="clear_tmpl"):
                del st.session_state["_ppt_template_key"]
                st.rerun()

    with st.expander("知识库（智能检索）", expanded=False):
        embedding_ready = kb.get_embedding_client() is not None
        if not embedding_ready:
            st.warning("未配置向量化服务（EMBEDDING_* 配置项：本地在 .env，线上在 Settings → Secrets），知识库暂不可用。")

        kb_files = st.file_uploader(
            "拖入或点击上传文档入库（txt/md/pdf/docx/pptx/xlsx，可多选）",
            type=UPLOAD_TYPES,
            key="kb_uploader",
            accept_multiple_files=True,
            disabled=not embedding_ready,
        )
        for kb_file in kb_files or []:
            done_key = f"kb_done_{kb_file.name}_{kb_file.size}"
            if done_key not in st.session_state:
                with st.spinner(f"正在将 {kb_file.name} 加入知识库…"):
                    try:
                        text = extract_text(kb_file)
                        n = knowledge_base.add_document(kb_file.name, text)
                        st.session_state[done_key] = (
                            "success",
                            f"{kb_file.name} 已入库（{n} 个分块）",
                        )
                    except Exception as e:
                        st.session_state[done_key] = ("error", f"{kb_file.name} 入库失败：{e}")
            status, msg = st.session_state[done_key]
            (st.success if status == "success" else st.error)(msg)

        doc_names = knowledge_base.doc_names()
        st.caption(f"知识库：{len(doc_names)} 个文档 / {len(knowledge_base.entries)} 个分块")
        if doc_names:
            for name in doc_names:
                if st.button(f"删除：{name}", key=f"del_{name}"):
                    knowledge_base.remove_document(name)
                    _get_kb.clear()
                    st.rerun()
            if st.button("清空知识库"):
                knowledge_base.clear()
                _get_kb.clear()
                st.rerun()

# ---------- 历史消息 ----------
for msg in messages:
    with st.chat_message(
        msg["role"],
        avatar=str(PENGUIN_IMG) if msg["role"] == "assistant" else None,
    ):
        st.markdown(msg["display"])
        if msg["role"] == "assistant" and msg.get("review"):
            with st.expander("🔍 自我审查结果", expanded=False):
                st.markdown(msg["review"])

# ---------- 操作条（紧贴输入框，界面更干净） ----------
last_reply = next(
    (m for m in reversed(messages) if m["role"] == "assistant" and not m.get("greeting")),
    None,
)
has_content = any(not m.get("greeting") for m in messages)
_ab0, _ab1, _ab2, _ab3, _ab4 = st.columns([1, 1, 1, 1, 1])
with _ab1:
    ppt_btn = st.button(
        "📊 生成 PPT",
        disabled=last_reply is None or client is None,
        help="将助手最新回复提炼为演示结构，套用模板生成 .pptx（含演讲备注）",
        use_container_width=True,
    )
with _ab2:
    doc_btn = st.button(
        "📤 导出 Word",
        disabled=not has_content,
        help="把当前话题完整导出为 Word 文档",
        use_container_width=True,
    )

# ---------- 生成处理 ----------
if ppt_btn and last_reply is not None and client is not None:
    with st.spinner("📊 正在分析内容并套用模板排版 PPT（约 10~30 秒）…"):
        src = last_reply.get("content") or last_reply.get("display") or ""
        outline = _ppt_outline(client, model_name, src) or _fallback_outline(src)
        st.session_state.ppt_bytes = _build_pptx(outline, cur["title"])
        st.session_state.ppt_autodl = True

if doc_btn and has_content:
    with st.spinner("📤 正在生成 Word 文档…"):
        st.session_state.docx_bytes = _build_docx(cur["title"], messages)
        st.session_state.docx_autodl = True

# ---------- 下载行 ----------
_dl_items = []
if st.session_state.get("ppt_bytes"):
    _dl_items.append(("ppt", st.session_state["ppt_bytes"],
                      f"{(cur['title'] or '演示文稿')[:20]}.pptx",
                      "application/vnd.openxmlformats-officedocument.presentationml.presentation"))
if st.session_state.get("docx_bytes"):
    _dl_items.append(("docx", st.session_state["docx_bytes"],
                      f"{(cur['title'] or '对话记录')[:20]}.docx",
                      "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))
if _dl_items:
    _d0, _d1, _d2, _d3, _d4 = st.columns([1, 1, 1, 1, 1])
    _dl_cols = [_d1, _d2]
    for col, (kind, data, fname, mime) in zip(_dl_cols, _dl_items):
        with col:
            st.download_button(
                "⬇️ 下载 PPT" if kind == "ppt" else "⬇️ 下载 Word",
                data=data, file_name=fname, mime=mime,
                key=f"{kind}_dl", use_container_width=True,
            )
    if st.session_state.pop("ppt_autodl", False):
        _autodownload(".pptx")
    if st.session_state.pop("docx_autodl", False):
        _autodownload(".docx")

# ---------- 输入与回复 ----------
submission = st.chat_input(
    "输入你的问题，📎可添加附件（悬浮📎查看支持格式）",
    accept_file="multiple",
    file_type=UPLOAD_TYPES,
)
if submission and ((submission.text or "").strip() or submission.files):
    user_input = (submission.text or "").strip()
    chat_files = list(submission.files or [])
    if not user_input and chat_files:
        user_input = "请结合我上传的附件内容进行回答。"

    if client is None:
        st.error("请先配置大模型密钥（OPENAI_API_KEY 或 EMBEDDING_API_KEY）：本地编辑 .env，线上在 Settings → Secrets。")
        st.stop()

    # 1) 本次附件文本（多文件拼接）
    upload_parts = []
    for f in chat_files:
        try:
            full = extract_text(f)
            part = full[:MAX_UPLOAD_CHARS]
            if len(full) > MAX_UPLOAD_CHARS:
                part += "\n（文件过长，已截断）"
            upload_parts.append(f"文件：{f.name}\n{part}")
        except Exception as e:
            upload_parts.append(f"文件：{f.name}（解析失败：{e}）")
    upload_text = "\n\n".join(upload_parts)

    # 2) 知识库检索
    kb_section = ""
    retrieval_hits = []
    if knowledge_base.entries and embedding_ready:
        try:
            hits = knowledge_base.search(user_input)
            retrieval_hits = hits
            if hits:
                kb_section = "【知识库检索资料】\n" + "\n\n".join(
                    f"[来源：{h['doc']}]\n{h['text']}" for h in hits
                )
        except Exception as e:
            st.warning(f"知识库检索失败：{e}")

    # 3) 组装发送给模型的用户消息
    parts = [
        p
        for p in [
            kb_section,
            f"【本次上传文件内容】\n{upload_text}" if upload_text else "",
            f"【用户问题】\n{user_input}",
        ]
        if p
    ]
    api_content = "\n\n".join(parts)

    # 4) 先把本次用户消息写入会话，再组装发送给模型的消息（开场白仅展示，不进上下文）
    display = user_input
    if chat_files:
        display += "\n\n📎 " + "、".join(f.name for f in chat_files)
    messages.append({"role": "user", "display": display, "content": api_content})

    messages_for_api = _truncate_history(
        [{"role": "system", "content": SYSTEM_PROMPT}]
        + [m for m in messages if not m.get("greeting")]
    )
    title_changed = False
    if sum(1 for m in messages if m["role"] == "user") == 1:  # 首条提问自动作为话题标题
        cur["title"] = user_input.replace("\n", " ")[:18]
        title_changed = True
    with st.chat_message("user"):
        st.markdown(display)

    with st.chat_message("assistant", avatar=str(PENGUIN_IMG)):
        try:
            stream = client.chat.completions.create(
                model=model_name,
                messages=messages_for_api,
                stream=True,
            )

            def _content_stream():
                for chunk in stream:
                    if not chunk.choices:  # 跳过空 choices 的统计块，避免越界
                        continue
                    content = chunk.choices[0].delta.content
                    if content:  # 跳过空内容块，避免界面显示 None
                        yield content

            reply = st.write_stream(_content_stream())
            assistant_msg = {"role": "assistant", "display": reply, "content": reply}
            messages.append(assistant_msg)

            # 5) 检索来源展示（可展开）
            if retrieval_hits:
                with st.expander(f"📚 检索来源（{len(retrieval_hits)} 条）", expanded=False):
                    for i, hit in enumerate(retrieval_hits, 1):
                        st.markdown(
                            f"**{i}. {hit['doc']}**  "
                            f"<small>相关度: {hit.get('score', 0):.2f}</small>",
                            unsafe_allow_html=True,
                        )
                        st.markdown(f"> {hit['text'][:200]}{'...' if len(hit['text']) > 200 else ''}")

            # 6) Agent 自我审查（可展开）
            if review_enabled and client is not None:
                with st.spinner("🔍 Agent 自我审查中…"):
                    review_result = _run_self_review(
                        client=client,
                        model_name=model_name,
                        query=user_input,
                        kb_context=kb_section,
                        upload_context=upload_text,
                        assistant_reply=reply,
                    )
                if review_result:
                    assistant_msg["review"] = review_result
                    verdict = "warn"
                    for line in review_result.split("\n"):
                        if "总体评级" in line:
                            if "pass" in line.lower():
                                verdict = "pass"
                            elif "fail" in line.lower():
                                verdict = "fail"
                            break
                    verdict_icon = {"pass": "✅", "warn": "⚠️", "fail": "❌"}.get(verdict, "⚠️")
                    verdict_label = {"pass": "通过", "warn": "注意", "fail": "不通过"}.get(verdict, "审查")
                    with st.expander(
                        f"{verdict_icon} 自我审查：{verdict_label}",
                        expanded=(verdict == "fail"),
                    ):
                        st.markdown(review_result)
                else:
                    st.caption("🔍 自我审查：审查服务暂不可用，已跳过。")
        except Exception as e:  # 网络/鉴权/模型错误，保留会话可重试
            messages.pop()
            st.error(f"调用失败：{e}")
    _save_topics(store)  # 话题内容本地持久化
    if title_changed:  # 侧边栏先于提问渲染，标题变化后刷新一次以同步显示
        st.rerun()