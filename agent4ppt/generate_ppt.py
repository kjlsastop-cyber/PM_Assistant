#!/usr/bin/env python3
"""
generate_ppt.py — generate-ppt skill for agent4ppt

Generates an editable PPTX file from a markdown file. The markdown
frontmatter specifies the PPTX template path and output filename.

Usage:
    python generate_ppt.py <markdown_file> [--output <output.pptx>]
"""

from __future__ import annotations

import argparse
import io
import os
import re
import sys
from pathlib import Path
from typing import Any, NamedTuple

# ---------------------------------------------------------------------------
# Dependency check
# ---------------------------------------------------------------------------

def _detect_lang() -> str:
    """Detect language from LANG environment variable."""
    env_lang = os.environ.get("LANG", "")
    lang_code = env_lang.split("_")[0].split(".")[0]
    return "ko" if lang_code == "ko" else "en"


def _try_auto_install(packages: list[str]) -> bool:
    """Attempt to auto-install missing packages via pip. Returns True on success."""
    import subprocess

    pip_cmd = [sys.executable, "-m", "pip", "install"] + packages
    strategies: list[tuple[str, list[str]]] = [
        ("pip install", []),
        ("pip install --user", ["--user"]),
    ]
    for label, extra_args in strategies:
        try:
            result = subprocess.run(
                pip_cmd + extra_args,
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode == 0:
                print(
                    f"[agent4ppt] Auto-installed via {label}: {', '.join(packages)}",
                    file=sys.stderr,
                )
                return True
        except (subprocess.TimeoutExpired, FileNotFoundError, OSError):
            continue
    return False


def _check_dependencies() -> None:
    """Check for required dependencies; attempt auto-install on failure."""
    lang = _detect_lang()
    missing: list[str] = []
    for pkg, name in [("pptx", "python-pptx"), ("yaml", "pyyaml"), ("markdown_it", "markdown-it-py")]:
        try:
            __import__(pkg)
        except ImportError:
            missing.append(name)

    if not missing:
        return

    # --- Attempt auto-install ---------------------------------------------------
    auto_install_env = os.environ.get("AGENT4PPT_AUTO_INSTALL", "1")
    if auto_install_env == "1":
        if lang == "ko":
            print(
                f"[agent4ppt] 누락된 의존성 자동 설치 시도 중: {', '.join(missing)}",
                file=sys.stderr,
            )
        else:
            print(
                f"[agent4ppt] Attempting auto-install of missing dependencies: {', '.join(missing)}",
                file=sys.stderr,
            )
        if _try_auto_install(missing):
            # Re-check after install
            still_missing = []
            for pkg, name in [("pptx", "python-pptx"), ("yaml", "pyyaml"), ("markdown_it", "markdown-it-py")]:
                try:
                    __import__(pkg)
                except ImportError:
                    still_missing.append(name)
            if not still_missing:
                return
            missing = still_missing

    # --- Provide fallback guidance ----------------------------------------------
    if lang == "ko":
        print(
            f"[agent4ppt] 누락된 의존성: {', '.join(missing)}\n"
            "설치 방법:\n"
            "  pip install " + " ".join(missing) + "\n"
            "또는 프로젝트 설치 스크립트 실행:\n"
            "  bash scripts/install.sh\n"
            "\n"
            "자동 설치를 비활성화하려면: AGENT4PPT_AUTO_INSTALL=0",
            file=sys.stderr,
        )
    else:
        print(
            f"[agent4ppt] Missing dependencies: {', '.join(missing)}\n"
            "Install with:\n"
            "  pip install " + " ".join(missing) + "\n"
            "Or run the project installer:\n"
            "  bash scripts/install.sh\n"
            "\n"
            "To disable auto-install: AGENT4PPT_AUTO_INSTALL=0",
            file=sys.stderr,
        )
    sys.exit(1)


_check_dependencies()

# ---------------------------------------------------------------------------
# Resolve the agent4ppt package path so this script can be run standalone
# (i.e. python skills/generate-ppt/generate_ppt.py)
# ---------------------------------------------------------------------------

_SKILL_DIR = Path(__file__).resolve().parent       # skills/generate-ppt/
_PROJECT_ROOT = _SKILL_DIR.parent.parent           # agent4ppt repo root

if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

# ---------------------------------------------------------------------------
# Imports
# ---------------------------------------------------------------------------

import yaml  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE  # noqa: E402
from pptx.chart.data import ChartData, XyChartData  # noqa: E402
from pptx.enum.shapes import PP_PLACEHOLDER  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

# Import shared library — available after sys.path setup above.
# Provides:
#   - parse_pptx_template / SlideTemplateInfo   (structural template info)
#   - _normalise_lang                           (POSIX locale → language code)
#   - parse_slide_section, apply_slide_content  (core generation logic)
#   - check_mismatches / MismatchReport         (template mismatch detection)
try:
    from agent4ppt.pptx_parser import parse_pptx_template, SlideTemplateInfo  # noqa: E402
    from agent4ppt.markdown_generator import _normalise_lang  # noqa: E402
    from agent4ppt.slide_generator import (  # noqa: E402
        parse_slide_section as _lib_parse_slide_section,
        apply_slide_content as _lib_apply_slide_content,
    )
    from agent4ppt.mismatch_detector import (  # noqa: E402
        check_mismatches as _lib_check_mismatches,
        MismatchReport,
    )
    from agent4ppt.image_validator import (  # noqa: E402
        validate_images_in_document as _lib_validate_images,
        SUPPORTED_IMAGE_EXTENSIONS as _LIB_SUPPORTED_IMG_EXTS,
    )
    _LIBRARY_AVAILABLE = True
except ImportError:
    _LIBRARY_AVAILABLE = False

# ---------------------------------------------------------------------------
# Language / message support
# ---------------------------------------------------------------------------

_MESSAGES: dict[str, dict[str, str]] = {
    "en": {
        "missing_template": "Error: template file not found: {path}",
        "missing_template_hint": (
            "Hint: The 'template' field in your markdown frontmatter points to '{path}'.\n"
            "  • Check that the file exists at this path.\n"
            "  • Relative paths are resolved from the markdown file's directory.\n"
            "  • To obtain a template, run: /parse-ppt-template <your_file.pptx>\n"
            "  • Or copy a .pptx file to the expected path and retry."
        ),
        "missing_md": "Error: markdown file not found: {path}",
        "missing_md_hint": (
            "Hint: Could not find the markdown content file at '{path}'.\n"
            "  • Check that the file path is correct.\n"
            "  • To generate a markdown template from a PPTX, run: /parse-ppt-template <your_file.pptx>"
        ),
        "no_template_field": "Error: 'template' field missing in frontmatter",
        "no_template_field_hint": (
            "Hint: Your markdown file must begin with YAML frontmatter that includes a 'template' field.\n"
            "  Example:\n"
            "    ---\n"
            "    template: ./my_template.pptx\n"
            "    fname: output.pptx\n"
            "    ---\n"
            "  • Run /parse-ppt-template <your_file.pptx> to generate a template with correct frontmatter."
        ),
        "layout_out_of_range": "Warning: layout index {idx} out of range (max {max}), using layout 0",
        "ph_not_found": "Warning: placeholder idx={idx} not found in slide layout {layout}, skipping",
        "img_not_found": "Warning: image file not found: {path}, skipping",
        "img_unsupported_format": (
            "Warning: unsupported image format '{ext}' for '{path}'. "
            "Supported: {supported}. Attempting insertion anyway."
        ),
        "img_insert_failed": "Warning: failed to insert image '{path}': {err}. Placeholder left empty.",
        "img_validation_header": "Image validation: {total} reference(s) found, {invalid} issue(s) detected.",
        "img_validation_warn": "  ⚠ {msg}",
        "invalid_chart": "Warning: invalid chart YAML in slide {slide}: {err}",
        "done": "PPTX generated → {path}",
        "loading_template": "Loading template: {path}",
        "loading_md": "Loading markdown: {path}",
        # Mismatch detector fallback messages (used when shared library unavailable)
        "mismatch_layout_range": (
            "Warning: Slide {slide}: layout index {idx} out of range "
            "(template has {count} layout(s), valid 0–{max}). "
            "Falling back to layout 0."
        ),
        "mismatch_ph_missing": (
            "Warning: Slide {slide}: placeholder idx={idx} (type={ph_type}) "
            "not found in layout {layout_idx}. Placeholder will be skipped."
        ),
        "mismatch_report_header": "[agent4ppt] ⚠  Template mismatch report ({count} warning(s))",
        "mismatch_report_continue": "[agent4ppt] Continuing generation with {count} warning(s) above.",
        "mismatch_report_none": "[agent4ppt] ✓  No template mismatches detected.",
    },
    "ko": {
        "missing_template": "오류: 템플릿 파일을 찾을 수 없습니다: {path}",
        "missing_template_hint": (
            "힌트: 마크다운 frontmatter의 'template' 필드가 '{path}'를 가리킵니다.\n"
            "  • 해당 경로에 파일이 존재하는지 확인하세요.\n"
            "  • 상대 경로는 마크다운 파일의 디렉토리 기준으로 해석됩니다.\n"
            "  • 템플릿을 얻으려면: /parse-ppt-template <your_file.pptx> 를 실행하세요.\n"
            "  • 또는 .pptx 파일을 해당 경로에 복사한 후 다시 시도하세요."
        ),
        "missing_md": "오류: 마크다운 파일을 찾을 수 없습니다: {path}",
        "missing_md_hint": (
            "힌트: '{path}'에서 마크다운 파일을 찾을 수 없습니다.\n"
            "  • 파일 경로가 올바른지 확인하세요.\n"
            "  • PPTX에서 마크다운 템플릿을 생성하려면: /parse-ppt-template <your_file.pptx> 를 실행하세요."
        ),
        "no_template_field": "오류: frontmatter에 'template' 필드가 없습니다",
        "no_template_field_hint": (
            "힌트: 마크다운 파일은 반드시 'template' 필드를 포함한 YAML frontmatter로 시작해야 합니다.\n"
            "  예시:\n"
            "    ---\n"
            "    template: ./my_template.pptx\n"
            "    fname: output.pptx\n"
            "    ---\n"
            "  • /parse-ppt-template <your_file.pptx> 를 실행하면 올바른 frontmatter가 포함된 템플릿을 생성합니다."
        ),
        "layout_out_of_range": "경고: 레이아웃 인덱스 {idx}가 범위를 벗어났습니다 (최대 {max}), 레이아웃 0 사용",
        "ph_not_found": "경고: 슬라이드 레이아웃 {layout}에서 idx={idx} 플레이스홀더를 찾을 수 없습니다, 건너뜀",
        "img_not_found": "경고: 이미지 파일을 찾을 수 없습니다: {path}, 건너뜀",
        "img_unsupported_format": (
            "경고: '{path}'의 이미지 형식 '{ext}'은(는) 지원되지 않습니다. "
            "지원 형식: {supported}. 삽입을 시도합니다."
        ),
        "img_insert_failed": "경고: 이미지 '{path}' 삽입 실패: {err}. 플레이스홀더를 비워 둡니다.",
        "img_validation_header": "이미지 검사: {total}개 참조 발견, {invalid}개 문제 감지.",
        "img_validation_warn": "  ⚠ {msg}",
        "invalid_chart": "경고: 슬라이드 {slide}에 유효하지 않은 차트 YAML: {err}",
        "done": "PPTX 생성 완료 → {path}",
        "loading_template": "템플릿 로드 중: {path}",
        "loading_md": "마크다운 로드 중: {path}",
        # Mismatch detector fallback messages (used when shared library unavailable)
        "mismatch_layout_range": (
            "경고: 슬라이드 {slide}: 레이아웃 인덱스 {idx}가 범위를 벗어났습니다 "
            "(템플릿에 {count}개 레이아웃, 유효 범위 0–{max}). 레이아웃 0으로 대체합니다."
        ),
        "mismatch_ph_missing": (
            "경고: 슬라이드 {slide}: 플레이스홀더 idx={idx} (type={ph_type})가 "
            "레이아웃 {layout_idx}에 없습니다. 건너뜁니다."
        ),
        "mismatch_report_header": "[agent4ppt] ⚠  템플릿 불일치 보고서 ({count}개 경고)",
        "mismatch_report_continue": "[agent4ppt] 위 {count}개 경고와 함께 생성을 계속합니다.",
        "mismatch_report_none": "[agent4ppt] ✓  템플릿 불일치 없음.",
    },
}


def _msg(lang: str, key: str, **kwargs) -> str:
    msgs = _MESSAGES.get(lang, _MESSAGES["en"])
    template = msgs.get(key, key)
    return template.format(**kwargs)


def _get_lang(env_override: str | None = None) -> str:
    """Determine language from argument, environment variable, or default.

    Delegates to _normalise_lang from the shared library when available,
    otherwise falls back to simple string parsing for standalone operation.
    """
    raw: str
    if env_override:
        raw = env_override
    else:
        raw = os.environ.get("LANG", "en")

    if _LIBRARY_AVAILABLE:
        return _normalise_lang(raw)
    # Fallback: simple parsing (covers en, ko, ko_KR.UTF-8, etc.)
    _supported = frozenset({"en", "ko"})
    code = raw.split("_")[0].split(".")[0].lower()
    return code if code in _supported else "en"

# ---------------------------------------------------------------------------
# Image support constants (standalone fallback — mirrors image_validator.py)
# ---------------------------------------------------------------------------

#: Supported image file extensions for local validation (same set as image_validator).
SUPPORTED_IMAGE_EXTENSIONS: frozenset[str] = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".svg"
})

_IMG_REF_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")


def _validate_images_local(
    markdown_text: str,
    md_dir: Path,
    lang: str,
) -> list[tuple[str, str]]:
    """Standalone image validation — returns a list of (path, warning_msg) pairs.

    Used when the shared agent4ppt library is unavailable.  Strips the YAML
    frontmatter block before scanning so frontmatter values don't produce false
    positives.  Remote URLs (http/https/ftp) are silently skipped.

    Args:
        markdown_text: Full raw markdown document text.
        md_dir:        Directory of the markdown file (for relative path resolution).
        lang:          Language code for formatting warnings (currently unused here).

    Returns:
        A list of ``(resolved_path_str, warning_message)`` tuples for each
        validation failure found.
    """
    _supported = SUPPORTED_IMAGE_EXTENSIONS
    body = markdown_text
    # Strip YAML frontmatter
    if body.lstrip().startswith("---"):
        try:
            first = body.index("---")
            rest = body[first + 3:]
            closing = rest.find("\n---")
            if closing != -1:
                body = rest[closing + 4:]
        except ValueError:
            pass

    issues: list[tuple[str, str]] = []
    for m in _IMG_REF_RE.finditer(body):
        alt = m.group(1)
        raw_path = m.group(2)
        if raw_path.startswith(("http://", "https://", "ftp://", "//")):
            continue
        img_path = Path(raw_path)
        resolved = img_path if img_path.is_absolute() else (md_dir / img_path)
        ext = resolved.suffix.lower()
        if ext not in _supported:
            supported_str = ", ".join(sorted(_supported))
            issues.append((
                str(resolved),
                f"Unsupported image format '{ext}' for '{raw_path}'. Supported: {supported_str}",
            ))
        if not resolved.exists():
            issues.append((
                str(resolved),
                f"Image file not found: '{resolved}' (referenced as '![{alt}]({raw_path})')",
            ))
    return issues


# ---------------------------------------------------------------------------
# Markdown parsing
# ---------------------------------------------------------------------------

_FRONTMATTER_RE = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)
_LAYOUT_RE = re.compile(r"^>\s*layout:\s*(\d+)", re.MULTILINE)
# NOTE: The annotation format generated by markdown_generator includes optional
# pos/size attributes after the type, e.g.:
#   <!-- ph:0 type:title pos:(0.5",0.5") size:8.0"×1.5" -->
# The .*? handles both the short form (no pos/size) and full form.
_PH_COMMENT_RE = re.compile(r"<!--\s*ph:(\d+)\s+type:(\w+).*?-->")
_GUIDE_COMMENT_RE = re.compile(r"<!--.*?-->", re.DOTALL)
_CHART_BLOCK_RE = re.compile(r"```chart\s*\n(.*?)```", re.DOTALL)
_IMAGE_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
_TABLE_ROW_RE = re.compile(r"^\|(.+)\|$")
_COLUMN_SEP = "|||"


def parse_frontmatter(text: str) -> tuple[dict[str, Any], str]:
    """Split YAML frontmatter from markdown body."""
    m = _FRONTMATTER_RE.match(text)
    if m:
        fm = yaml.safe_load(m.group(1)) or {}
        body = text[m.end():]
        return fm, body
    return {}, text


def split_slides(body: str) -> list[str]:
    """Split body on --- slide separators."""
    # Split on lines that are exactly '---'
    parts = re.split(r"\n---\s*\n", body)
    return [p.strip() for p in parts if p.strip()]


def parse_slide_section(section: str) -> dict[str, Any]:
    """Parse a single slide section into structured data.

    Delegates to ``agent4ppt.slide_generator.parse_slide_section`` when the
    shared library is available, otherwise falls back to a local
    implementation using the same fixed regex that handles both the short
    annotation form (``<!-- ph:N type:T -->``) and the full form generated by
    ``markdown_generator`` (``<!-- ph:N type:T pos:... size:... -->``).
    """
    if _LIBRARY_AVAILABLE:
        return _lib_parse_slide_section(section)

    # ── Standalone fallback ──────────────────────────────────────────────────
    slide: dict[str, Any] = {
        "layout": 0,
        "placeholders": [],  # list of {idx, type, content}
    }

    m = _LAYOUT_RE.search(section)
    if m:
        slide["layout"] = int(m.group(1))

    # Split on <!-- ph:N type:T ... --> markers.
    # _PH_COMMENT_RE captures (idx, type) and the .*? handles optional attrs.
    parts = _PH_COMMENT_RE.split(section)
    # parts: [pre, idx1, type1, content1, idx2, type2, content2, ...]
    i = 1
    while i + 2 < len(parts):
        ph_idx = int(parts[i])
        ph_type = parts[i + 1]
        content = parts[i + 2].strip()
        # Remove remaining guide comments from content
        content = _GUIDE_COMMENT_RE.sub("", content).strip()
        slide["placeholders"].append({
            "idx": ph_idx,
            "type": ph_type,
            "content": content,
        })
        i += 3

    return slide


# ---------------------------------------------------------------------------
# Loaded data structures
# ---------------------------------------------------------------------------

class MarkdownDocument(NamedTuple):
    """Parsed result of reading and splitting a markdown content file."""
    frontmatter: dict[str, Any]   # YAML frontmatter key-value pairs
    body: str                      # Markdown body text (after frontmatter)
    slide_sections: list[str]      # Raw slide sections split on ---
    source_path: Path              # Absolute path to the source .md file
    lang: str                      # Resolved language code (en/ko)
    raw_content: str = ""          # Full original file text (for version hashing)


class LoadedTemplate(NamedTuple):
    """Result of loading a PPTX template for PPTX generation."""
    presentation: Any              # pptx.Presentation object ready for slide generation
    template_path: Path            # Absolute path to the resolved .pptx template
    template_info: Any | None      # SlideTemplateInfo from shared library (or None)


def load_markdown_and_template(
    md_path: Path,
    lang: str,
) -> tuple[MarkdownDocument, LoadedTemplate]:
    """Load a markdown content file and its referenced PPTX template.

    This is the primary entry point for the /generate-ppt skill.  It:

    1. Reads the markdown file and parses YAML frontmatter.
    2. Extracts the ``template`` frontmatter field and resolves it relative
       to the markdown file's directory.
    3. Loads the PPTX template using python-pptx ``Presentation``.
    4. If the shared ``agent4ppt`` library is available, also calls
       ``parse_pptx_template()`` to obtain structured layout/placeholder
       metadata (``SlideTemplateInfo``) useful for validation and richer
       error messages.
    5. Splits the markdown body into per-slide sections ready for the
       generation phase.

    Args:
        md_path: Path to the markdown file to read.
        lang: Language code for error messages (``"en"`` or ``"ko"``).

    Returns:
        A ``(MarkdownDocument, LoadedTemplate)`` tuple.  Callers can
        immediately pass these to ``generate_ppt()`` or inspect them for
        validation.

    Raises:
        SystemExit(1) on any unrecoverable error (missing file, missing
        frontmatter field, etc.) after printing a human-readable message.
    """
    # ── 1. Read markdown -------------------------------------------------
    text = md_path.read_text(encoding="utf-8")
    fm, body = parse_frontmatter(text)

    # Language from frontmatter overrides CLI / env detection
    lang = _get_lang(fm.get("lang") or lang)

    # ── 2. Resolve PPTX template path -----------------------------------
    template_field = fm.get("template")
    if not template_field:
        print(f"[agent4ppt] {_msg(lang, 'no_template_field')}", file=sys.stderr)
        print(_msg(lang, "no_template_field_hint"), file=sys.stderr)
        sys.exit(1)

    template_path = Path(template_field)
    if not template_path.is_absolute():
        template_path = md_path.parent / template_path

    if not template_path.exists():
        print(
            f"[agent4ppt] {_msg(lang, 'missing_template', path=template_path)}",
            file=sys.stderr,
        )
        print(
            _msg(lang, "missing_template_hint", path=template_field),
            file=sys.stderr,
        )
        sys.exit(1)

    # ── 3. Load Presentation (python-pptx) --------------------------------
    prs = Presentation(str(template_path))

    # ── 4. Optionally parse structural template info (shared library) ----
    template_info: Any | None = None
    if _LIBRARY_AVAILABLE:
        try:
            template_info = parse_pptx_template(str(template_path))
        except Exception:
            # Non-fatal: structural info is used for validation only
            template_info = None

    # ── 5. Split body into slide sections --------------------------------
    slide_sections = split_slides(body)

    md_doc = MarkdownDocument(
        frontmatter=fm,
        body=body,
        slide_sections=slide_sections,
        source_path=md_path.resolve(),
        lang=lang,
        raw_content=text,
    )
    loaded_tpl = LoadedTemplate(
        presentation=prs,
        template_path=template_path.resolve(),
        template_info=template_info,
    )
    return md_doc, loaded_tpl


# ---------------------------------------------------------------------------
# Content renderers
# ---------------------------------------------------------------------------

def _render_text_to_tf(text_frame, content: str) -> None:
    """Render markdown text into a TextFrame, supporting bullets/bold/italic."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    lines = content.splitlines()
    first_para = True
    for line in lines:
        line = line.rstrip()
        # Skip empty lines
        if not line:
            continue

        # Detect bullet level
        bullet_match = re.match(r"^(\s*)([-*+]|\d+\.)\s+(.*)", line)
        if bullet_match:
            indent = len(bullet_match.group(1))
            level = min(indent // 2, 8)
            text = bullet_match.group(3)
        else:
            level = 0
            text = line

        if first_para:
            para = text_frame.paragraphs[0]
            first_para = False
        else:
            para = text_frame.add_paragraph()

        para.level = level
        _render_inline_markdown(para, text)


def _render_inline_markdown(para, text: str) -> None:
    """Render inline markdown (bold, italic, links) into a paragraph."""
    from pptx.util import Pt

    # Tokenize: bold, italic, links, plain
    pattern = re.compile(
        r"(\*\*(.+?)\*\*)"       # bold
        r"|(\*(.+?)\*)"          # italic
        r"|(\[([^\]]+)\]\(([^)]+)\))"  # link
    )

    pos = 0
    for m in pattern.finditer(text):
        # Plain text before match
        if m.start() > pos:
            run = para.add_run()
            run.text = text[pos:m.start()]

        if m.group(1):  # bold
            run = para.add_run()
            run.text = m.group(2)
            run.font.bold = True
        elif m.group(3):  # italic
            run = para.add_run()
            run.text = m.group(4)
            run.font.italic = True
        elif m.group(5):  # link
            run = para.add_run()
            run.text = m.group(6)
            # Hyperlinks require relationship; just underline as visual hint
            run.font.underline = True

        pos = m.end()

    # Remaining plain text
    if pos < len(text):
        run = para.add_run()
        run.text = text[pos:]


def _parse_table_alignment_local(separator_cells: list) -> list:
    """Parse GFM column alignment from separator row cells (local fallback).

    Returns list of PP_ALIGN values, one per cell.
    """
    alignments = []
    for cell in separator_cells:
        cell = cell.strip()
        starts_colon = cell.startswith(":")
        ends_colon = cell.endswith(":")
        if starts_colon and ends_colon:
            alignments.append(PP_ALIGN.CENTER)
        elif ends_colon:
            alignments.append(PP_ALIGN.RIGHT)
        else:
            alignments.append(PP_ALIGN.LEFT)
    return alignments


def _insert_table(slide, placeholder, content: str, lang: str, slide_num: int) -> None:
    """Replace a placeholder with a python-pptx table parsed from markdown.

    Supports:
    * GFM column alignment from separator rows (``:---``, ``---:``, ``:---:``)
    * Inline markdown (bold, italic) in table cells
    * Header row bolding
    """
    rows_data = []
    col_alignments = []

    for line in content.splitlines():
        line = line.strip()
        if not line:
            continue
        if not _TABLE_ROW_RE.match(line):
            continue

        cells = [c.strip() for c in line.strip("|").split("|")]

        # Detect separator row: all cells match [-: ]+
        if cells and all(re.match(r"^[-: ]+$", c) for c in cells if c):
            col_alignments = _parse_table_alignment_local(cells)
            continue

        rows_data.append(cells)

    if not rows_data:
        return

    num_rows = len(rows_data)
    num_cols = max(len(r) for r in rows_data)

    # Pad alignment list
    while len(col_alignments) < num_cols:
        col_alignments.append(PP_ALIGN.LEFT)

    # Position/size from placeholder
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # Remove placeholder shape first
    sp = placeholder._element
    sp.getparent().remove(sp)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    for r_idx, row in enumerate(rows_data):
        for c_idx in range(num_cols):
            cell_text = row[c_idx] if c_idx < len(row) else ""
            cell = table.cell(r_idx, c_idx)

            tf = cell.text_frame
            tf.clear()
            para = tf.paragraphs[0]

            if c_idx < len(col_alignments):
                para.alignment = col_alignments[c_idx]

            _render_inline_markdown(para, cell_text)

            if r_idx == 0:
                for run in para.runs:
                    run.font.bold = True


_LOCAL_CHART_TYPE_MAP: dict[str, "XL_CHART_TYPE"] = {
    # Bar (horizontal)
    "bar":                       XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_clustered":             XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked":               XL_CHART_TYPE.BAR_STACKED,
    "bar_stacked_100":           XL_CHART_TYPE.BAR_STACKED_100,
    "bar_percent":               XL_CHART_TYPE.BAR_STACKED_100,
    # Column (vertical)
    "column":                    XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_clustered":          XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked":            XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100":        XL_CHART_TYPE.COLUMN_STACKED_100,
    "column_percent":            XL_CHART_TYPE.COLUMN_STACKED_100,
    # Line
    "line":                      XL_CHART_TYPE.LINE,
    "line_markers":              XL_CHART_TYPE.LINE_MARKERS,
    "line_with_markers":         XL_CHART_TYPE.LINE_MARKERS,
    "line_stacked":              XL_CHART_TYPE.LINE_STACKED,
    "line_stacked_100":          XL_CHART_TYPE.LINE_STACKED_100,
    "line_stacked_markers":      XL_CHART_TYPE.LINE_MARKERS_STACKED,
    # Pie / Doughnut
    "pie":                       XL_CHART_TYPE.PIE,
    "pie_exploded":              XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut":                  XL_CHART_TYPE.DOUGHNUT,
    "doughnut_exploded":         XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    # Area
    "area":                      XL_CHART_TYPE.AREA,
    "area_stacked":              XL_CHART_TYPE.AREA_STACKED,
    "area_stacked_100":          XL_CHART_TYPE.AREA_STACKED_100,
    "area_percent":              XL_CHART_TYPE.AREA_STACKED_100,
    # Scatter / XY
    "scatter":                   XL_CHART_TYPE.XY_SCATTER,
    "xy_scatter":                XL_CHART_TYPE.XY_SCATTER,
    "scatter_lines":             XL_CHART_TYPE.XY_SCATTER_LINES,
    "scatter_smooth":            XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    "scatter_lines_no_markers":  XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    "scatter_smooth_no_markers": XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
    # Radar
    "radar":                     XL_CHART_TYPE.RADAR,
    "radar_filled":              XL_CHART_TYPE.RADAR_FILLED,
    "radar_markers":             XL_CHART_TYPE.RADAR_MARKERS,
}

_LOCAL_SCATTER_TYPES = frozenset({
    "scatter", "xy_scatter", "scatter_lines", "scatter_smooth",
    "scatter_lines_no_markers", "scatter_smooth_no_markers",
})
_LOCAL_NO_AXIS_TYPES = frozenset({
    "pie", "pie_exploded", "doughnut", "doughnut_exploded",
    "radar", "radar_filled", "radar_markers",
})


def _insert_chart(slide, placeholder, content: str, lang: str, slide_num: int) -> None:
    """Replace a placeholder with a python-pptx native chart.

    Supports bar, column, line (with optional markers), pie, doughnut, area,
    scatter/XY and radar chart types. Chart title, legend, data labels and
    axis titles are configurable from the YAML block.

    See agent4ppt.slide_generator._insert_chart for the full YAML reference.
    """
    chart_match = _CHART_BLOCK_RE.search(content)
    if not chart_match:
        return

    try:
        chart_def = yaml.safe_load(chart_match.group(1))
    except yaml.YAMLError as exc:
        print(f"[agent4ppt] {_msg(lang, 'invalid_chart', slide=slide_num, err=exc)}")
        return

    if not isinstance(chart_def, dict):
        print(
            f"[agent4ppt] Warning: chart YAML in slide {slide_num + 1} must be a mapping, skipping",
            file=sys.stderr,
        )
        return

    chart_type_str = str(chart_def.get("type", "bar")).lower().replace("-", "_")
    xl_type = _LOCAL_CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.BAR_CLUSTERED)
    is_scatter = chart_type_str in _LOCAL_SCATTER_TYPES
    has_axes = chart_type_str not in _LOCAL_NO_AXIS_TYPES

    categories = chart_def.get("categories", [])
    series_list = chart_def.get("series", [])

    if not series_list:
        return

    # Build chart data
    if is_scatter:
        chart_data = XyChartData()
        for series in series_list:
            series_obj = chart_data.add_series(series.get("name", ""))
            x_vals = series.get("x_values", series.get("values", []))
            y_vals = series.get("y_values", [])
            if y_vals:
                for x_val, y_val in zip(x_vals, y_vals):
                    series_obj.add_data_point(float(x_val), float(y_val))
            else:
                for i, y_val in enumerate(x_vals):
                    series_obj.add_data_point(float(i), float(y_val))
    else:
        chart_data = ChartData()
        chart_data.categories = [str(c) for c in categories]
        for series in series_list:
            chart_data.add_series(
                series.get("name", ""),
                [float(v) for v in series.get("values", [])],
            )

    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    sp = placeholder._element
    sp.getparent().remove(sp)

    graphic_frame = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart

    # Chart title
    title_text = str(chart_def.get("title", "")).strip()
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text

    # Legend — pie/doughnut always show legend by default (matches add_pie_chart behaviour)
    _PIE_TYPES_LOCAL = frozenset({"pie", "pie_exploded", "doughnut", "doughnut_exploded"})
    if chart_type_str in _PIE_TYPES_LOCAL:
        default_legend = True
    else:
        default_legend = len(series_list) > 1
    show_legend = chart_def.get("legend", default_legend)
    chart.has_legend = bool(show_legend)

    # Data labels
    if chart_def.get("data_labels", False):
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
        except Exception:
            pass

    # Axis titles
    if has_axes:
        x_axis_title = str(chart_def.get("x_axis", "")).strip()
        y_axis_title = str(chart_def.get("y_axis", "")).strip()

        if x_axis_title:
            try:
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = x_axis_title
            except Exception:
                pass

        if y_axis_title:
            try:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = y_axis_title
            except Exception:
                pass


def _insert_image(slide, placeholder, image_path: str, lang: str, md_dir: Path) -> None:
    """Replace a placeholder with an image, preserving aspect ratio.

    Scales the image to fit within the placeholder bounds while maintaining
    the original aspect ratio. The result is centred in the placeholder area.
    Falls back to stretching if aspect-ratio handling fails.

    Error handling:
    * File not found → warning printed, placeholder left as-is.
    * Unsupported extension → warning printed, insertion attempted anyway
      (python-pptx may still handle it; better a stretched image than nothing).
    * Insert failure (corrupt file, truly unsupported format) → warning printed,
      placeholder removed but no picture added.
    """
    img_path = Path(image_path)
    if not img_path.is_absolute():
        img_path = md_dir / img_path

    if not img_path.exists():
        print(
            f"[agent4ppt] {_msg(lang, 'img_not_found', path=img_path)}",
            file=sys.stderr,
        )
        return

    # --- Extension check (warn but proceed — python-pptx may still support it)
    _supported = _LIB_SUPPORTED_IMG_EXTS if _LIBRARY_AVAILABLE else SUPPORTED_IMAGE_EXTENSIONS
    ext = img_path.suffix.lower()
    if ext not in _supported:
        supported_str = ", ".join(sorted(_supported))
        print(
            f"[agent4ppt] {_msg(lang, 'img_unsupported_format', ext=ext, path=img_path, supported=supported_str)}",
            file=sys.stderr,
        )

    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    sp = placeholder._element
    sp.getparent().remove(sp)

    try:
        pic = slide.shapes.add_picture(str(img_path), left, top, width, None)
        if pic.height > height:
            pic._element.getparent().remove(pic._element)
            pic = slide.shapes.add_picture(str(img_path), left, top, None, height)
            pic.left = left + (width - pic.width) // 2
        else:
            pic.top = top + (height - pic.height) // 2
    except Exception as exc:
        # Fallback: try stretching to fill placeholder bounds
        try:
            slide.shapes.add_picture(str(img_path), left, top, width, height)
        except Exception as exc2:
            # Both strategies failed — report the error so the user knows
            print(
                f"[agent4ppt] {_msg(lang, 'img_insert_failed', path=img_path, err=exc2)}",
                file=sys.stderr,
            )


def _split_columns(content: str) -> list[str]:
    """Split content by ||| multi-column separator."""
    return [c.strip() for c in content.split(_COLUMN_SEP)]


def _insert_multicolumn_local(
    slide,
    placeholder,
    content: str,
    lang: str,
    slide_num: int,
) -> None:
    """Render multi-column content as side-by-side text boxes (local fallback).

    Mirrors ``agent4ppt.slide_generator._insert_multicolumn`` for use when
    the shared library is not available.
    """
    columns = [c.strip() for c in content.split(_COLUMN_SEP)]
    columns = [c for c in columns if c][:4]
    if not columns:
        return

    n_cols = len(columns)
    left = placeholder.left
    top = placeholder.top
    total_width = placeholder.width
    height = placeholder.height

    sp = placeholder._element
    sp.getparent().remove(sp)

    gap = max(int(total_width * 0.05), 45720)
    col_width = (total_width - gap * (n_cols - 1)) // n_cols

    for col_idx, col_content in enumerate(columns):
        col_left = left + col_idx * (col_width + gap)
        txBox = slide.shapes.add_textbox(col_left, top, col_width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        _render_text_to_tf(tf, col_content)


def _apply_slide_content_local(
    slide,
    slide_data: dict[str, Any],
    md_dir: Path,
    lang: str,
    slide_num: int,
) -> dict[str, int]:
    """Standalone fallback for apply_slide_content (no shared library).

    Mirrors ``agent4ppt.slide_generator.apply_slide_content`` exactly so that
    the generate-ppt skill works correctly even when the agent4ppt package is
    not on sys.path (e.g. during standalone testing).

    Rich content dispatch:
    * chart block → native chart
    * picture type → image (aspect-ratio-preserving)
    * table/object/body + markdown table rows → table with alignment + inline md
    * object/body + image → aspect-ratio-preserving image
    * body/object/text + ||| → side-by-side text boxes
    * other types + ||| → first column only (safe fallback)
    * title-like → strip heading markers
    * otherwise → render_text_to_tf

    Returns:
        A dict with ``"total"``, ``"filled"``, ``"skipped"`` placeholder counts.
    """
    ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    total_phs = len(slide_data["placeholders"])
    filled_phs = 0
    skipped_phs = 0

    for ph_data in slide_data["placeholders"]:
        idx = ph_data["idx"]
        ph_type = ph_data["type"]
        content = ph_data["content"]

        if idx not in ph_map:
            skipped_phs += 1
            print(f"[agent4ppt] {_msg(lang, 'ph_not_found', idx=idx, layout=slide_data['layout'])}")
            continue

        ph = ph_map[idx]

        # ── Chart ──────────────────────────────────────────────────────────────
        if _CHART_BLOCK_RE.search(content):
            _insert_chart(slide, ph, content, lang, slide_num)
            filled_phs += 1
            continue

        # ── Picture placeholder ────────────────────────────────────────────────
        if ph_type == "picture":
            img_match = _IMAGE_RE.search(content)
            if img_match:
                _insert_image(slide, ph, img_match.group(2), lang, md_dir)
                filled_phs += 1
            continue

        # ── Multi-column layout (body / object / text) ─────────────────────────
        # IMPORTANT: check BEFORE table detection because "|||" matches
        # _TABLE_ROW_RE (starts and ends with "|") and would mis-route.
        if _COLUMN_SEP in content:
            if ph_type in ("body", "object", "text"):
                _insert_multicolumn_local(slide, ph, content, lang, slide_num)
                filled_phs += 1
                continue
            else:
                # Safe fallback: use first column and fall through to text render
                content = _split_columns(content)[0]

        # ── Markdown table (table / object / body) ─────────────────────────────
        table_lines = [ln for ln in content.splitlines() if _TABLE_ROW_RE.match(ln.strip())]
        if table_lines and ph_type in ("table", "object", "body"):
            _insert_table(slide, ph, content, lang, slide_num)
            filled_phs += 1
            continue

        # ── Image in object / body placeholder ────────────────────────────────
        img_match = _IMAGE_RE.search(content)
        if img_match and ph_type in ("object", "body"):
            _insert_image(slide, ph, img_match.group(2), lang, md_dir)
            filled_phs += 1
            continue

        # ── Strip heading markers for title-like placeholders ──────────────────
        if ph_type in ("title", "subtitle", "center_title"):
            content = re.sub(r"^#+\s*", "", content).strip()

        # ── Render text content ────────────────────────────────────────────────
        try:
            tf = ph.text_frame
            tf.clear()
            _render_text_to_tf(tf, content)
            filled_phs += 1
        except Exception:
            try:
                ph.text = content
                filled_phs += 1
            except Exception:
                skipped_phs += 1

    return {"total": total_phs, "filled": filled_phs, "skipped": skipped_phs}

# ---------------------------------------------------------------------------
# Standalone mismatch detection fallback
# ---------------------------------------------------------------------------

def _check_mismatches_local(
    slide_sections: list[str],
    prs,
    lang: str,
) -> list[str]:
    """Lightweight standalone mismatch detector (no shared library required).

    Checks for:
    * LAYOUT_OUT_OF_RANGE: layout index >= number of layouts in the template.
    * PLACEHOLDER_NOT_IN_LAYOUT: ph idx not in the layout's placeholder list.

    Returns a list of human-readable warning strings.  Generation continues
    regardless; callers are responsible for printing the warnings.
    """
    warnings_found: list[str] = []
    num_layouts = len(prs.slide_layouts)

    for slide_num_0, section in enumerate(slide_sections):
        slide_num = slide_num_0 + 1

        layout_match = _LAYOUT_RE.search(section)
        layout_idx = int(layout_match.group(1)) if layout_match else 0

        if layout_idx >= num_layouts:
            warnings_found.append(
                _msg(
                    lang,
                    "mismatch_layout_range",
                    slide=slide_num,
                    idx=layout_idx,
                    count=num_layouts,
                    max=max(num_layouts - 1, 0),
                )
            )
            layout_idx = 0

        # Build the idx set for the resolved layout
        layout = prs.slide_layouts[layout_idx]
        layout_ph_idxs: set[int] = {
            ph.placeholder_format.idx for ph in layout.placeholders
        }

        for m in _PH_COMMENT_RE.finditer(section):
            ph_idx = int(m.group(1))
            ph_type = m.group(2)
            if ph_idx not in layout_ph_idxs:
                warnings_found.append(
                    _msg(
                        lang,
                        "mismatch_ph_missing",
                        slide=slide_num,
                        idx=ph_idx,
                        ph_type=ph_type,
                        layout_idx=layout_idx,
                    )
                )

    return warnings_found


def _print_local_mismatch_report(warnings_found: list[str], lang: str) -> None:
    """Print the local (standalone) mismatch report to stderr."""
    if not warnings_found:
        print(_msg(lang, "mismatch_report_none"), file=sys.stderr)
        return
    print(
        _msg(lang, "mismatch_report_header", count=len(warnings_found)),
        file=sys.stderr,
    )
    for w in warnings_found:
        print(f"  ⚠ {w}", file=sys.stderr)
    print(
        _msg(lang, "mismatch_report_continue", count=len(warnings_found)),
        file=sys.stderr,
    )


def _print_partial_fill_summary_local(
    total: int, filled: int, skipped: int, lang: str
) -> None:
    """Print a partial-fill summary using standalone messages (no shared library)."""
    if skipped <= 0:
        return
    if lang == "ko":
        msg = (
            f"[agent4ppt] ℹ  부분 채움: {total}개 중 {filled}개 플레이스홀더 채움; "
            f"불일치로 {skipped}개 건너뜀."
        )
    else:
        msg = (
            f"[agent4ppt] ℹ  Partial fill: {filled} of {total} placeholder(s) filled; "
            f"{skipped} skipped due to mismatches."
        )
    print(msg, file=sys.stderr)


# ---------------------------------------------------------------------------
# Image validation helper
# ---------------------------------------------------------------------------

def _run_image_validation(markdown_text: str, md_dir: Path, lang: str) -> None:
    """Run a pre-generation image validation pass and print consolidated warnings.

    Uses the shared ``agent4ppt.image_validator`` when available, otherwise
    falls back to the local :func:`_validate_images_local` implementation.

    Issues found are printed to *stderr* as warnings.  Generation is never
    aborted; the warnings are informational only.

    Args:
        markdown_text: Full raw markdown document text.
        md_dir:        Directory of the markdown file (for path resolution).
        lang:          Language code for user-facing messages.
    """
    if _LIBRARY_AVAILABLE:
        try:
            results = _lib_validate_images(markdown_text, md_dir, emit_warnings=False, lang=lang)
            invalid = [r for r in results if not r.is_valid]
            if results:
                if invalid:
                    print(
                        f"[agent4ppt] {_msg(lang, 'img_validation_header', total=len(results), invalid=len(invalid))}",
                        file=sys.stderr,
                    )
                    for r in invalid:
                        for w in r.warnings:
                            print(
                                f"[agent4ppt] {_msg(lang, 'img_validation_warn', msg=w)}",
                                file=sys.stderr,
                            )
            return
        except Exception:
            pass  # Fall through to local validation on unexpected library error

    # Standalone fallback
    issues = _validate_images_local(markdown_text, md_dir, lang)
    if issues:
        # Deduplicate: count unique path+message pairs
        print(
            f"[agent4ppt] {_msg(lang, 'img_validation_header', total=len(issues), invalid=len(issues))}",
            file=sys.stderr,
        )
        for _path, msg in issues:
            print(
                f"[agent4ppt] {_msg(lang, 'img_validation_warn', msg=msg)}",
                file=sys.stderr,
            )


# ---------------------------------------------------------------------------
# Template slide management
# ---------------------------------------------------------------------------

def _clear_template_slides(prs) -> None:
    """Remove all existing slides from a Presentation object.

    The PPTX template may contain sample slides.  We remove them so the
    output contains only the slides defined in the markdown file.  The slide
    masters and layouts (which carry the visual styles) are preserved.

    Args:
        prs: A ``pptx.Presentation`` object whose existing slides to remove.
    """
    xml_slides = prs.slides._sldIdLst
    # Iterate in reverse to avoid index shifting
    for i in range(len(xml_slides) - 1, -1, -1):
        rId = xml_slides[i].rId
        prs.part.drop_rel(rId)
        del xml_slides[i]


# ---------------------------------------------------------------------------
# Main generation logic
# ---------------------------------------------------------------------------

def generate_ppt(md_path: Path, output_path: Path | None, lang: str, *, dry_run: bool = False) -> Path:
    """Generate PPTX from markdown file.

    Delegates loading to :func:`load_markdown_and_template` which reads the
    markdown, parses frontmatter, resolves the template reference, and loads
    both the presentation object and (when available) the structured template
    info from the shared agent4ppt library.

    When *dry_run* is ``True`` the entire pipeline runs (parsing, validation,
    slide construction in memory) but the final ``.save()`` call is skipped.
    This allows the research-and-present pipeline to verify content/template
    compatibility without producing a PPTX file on disk.
    """
    md_doc, loaded_tpl = load_markdown_and_template(md_path, lang)

    # Use resolved language from loaded markdown document
    lang = md_doc.lang
    fm = md_doc.frontmatter
    prs = loaded_tpl.presentation

    # Determine output path
    if output_path is None:
        fname = fm.get("fname")
        if fname:
            output_path = md_path.parent / fname
        else:
            output_path = md_path.with_suffix(".pptx")

    # Remove all existing slides from the template so the output contains
    # only the slides defined in the markdown (template is used for styles
    # and layouts only, not for content).
    _clear_template_slides(prs)

    # ── Mismatch detection (pre-generation validation pass) ──────────────────
    # Run before the generation loop so users see a consolidated warning report
    # rather than one message per skipped placeholder.  Generation is never
    # halted by mismatches; warnings are printed to stderr and processing
    # continues with safe fallbacks (e.g. layout 0 for out-of-range layouts).
    mismatch_report = None  # may be set below for partial-fill summary
    if _LIBRARY_AVAILABLE and loaded_tpl.template_info is not None:
        mismatch_report = _lib_check_mismatches(
            md_doc.slide_sections,
            loaded_tpl.template_info,
            lang=lang,
            markdown_content=md_doc.raw_content or None,
            template_path=str(loaded_tpl.template_path),
        )
        mismatch_report.print_report()
    else:
        # Standalone fallback: lightweight checks using only python-pptx layout info
        local_warnings = _check_mismatches_local(md_doc.slide_sections, prs, lang)
        _print_local_mismatch_report(local_warnings, lang)

    # ── Image validation (pre-generation pass) ───────────────────────────────
    # Scan all image references in the markdown for missing files and
    # unsupported extensions *before* starting slide construction so users
    # see a consolidated report rather than one warning per missing image.
    # Generation is never aborted by image issues; missing images simply leave
    # the placeholder empty while a clear warning is printed.
    _run_image_validation(md_doc.raw_content or "", md_path.parent, lang)

    num_layouts = len(prs.slide_layouts)

    # ── Partial-fill tracking across all slides ──────────────────────────────
    agg_total = 0
    agg_filled = 0
    agg_skipped = 0

    for slide_num, section in enumerate(md_doc.slide_sections):
        slide_data = parse_slide_section(section)
        layout_idx = slide_data["layout"]

        if layout_idx >= num_layouts:
            print(f"[agent4ppt] {_msg(lang, 'layout_out_of_range', idx=layout_idx, max=num_layouts - 1)}")
            layout_idx = 0

        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Delegate per-slide content application to the shared library
        # (or local fallback when the library is unavailable)
        if _LIBRARY_AVAILABLE:
            stats = _lib_apply_slide_content(
                slide, slide_data, md_path.parent, lang, slide_num
            )
        else:
            stats = _apply_slide_content_local(
                slide, slide_data, md_path.parent, lang, slide_num
            )

        # Accumulate partial-fill stats (if returned)
        if isinstance(stats, dict):
            agg_total += stats.get("total", 0)
            agg_filled += stats.get("filled", 0)
            agg_skipped += stats.get("skipped", 0)

    # ── Print partial-fill summary if any placeholders were skipped ───────────
    if agg_skipped > 0:
        if _LIBRARY_AVAILABLE and hasattr(mismatch_report, "print_partial_fill_summary"):
            mismatch_report.print_partial_fill_summary(agg_total, agg_filled)
        else:
            _print_partial_fill_summary_local(agg_total, agg_filled, agg_skipped, lang)

    if dry_run:
        print(f"[agent4ppt] Dry-run OK — {len(md_doc.slide_sections)} slide(s) validated, no file written")
        return output_path
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[agent4ppt] {_msg(lang, 'done', path=output_path)}")
    return output_path


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="generate-ppt",
        description="Generate a PPTX from a markdown file.",
    )
    parser.add_argument("markdown_file", help="Path to the markdown content file")
    parser.add_argument(
        "--output", "-o", default=None,
        help="Output PPTX file path (default: frontmatter fname or <input>.pptx)",
    )
    parser.add_argument(
        "--lang", "-l", default=None, choices=["ko", "en"],
        help="Language for messages: ko | en (default: $LANG or en)",
    )
    parser.add_argument(
        "--dry-run", action="store_true", default=False,
        help="Validate inputs and build slides in memory without writing the PPTX file",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    md_path = Path(args.markdown_file)
    if not md_path.exists():
        lang = _get_lang(args.lang)
        print(f"[agent4ppt] {_msg(lang, 'missing_md', path=md_path)}", file=sys.stderr)
        print(_msg(lang, "missing_md_hint", path=md_path), file=sys.stderr)
        return 1

    output_path = Path(args.output) if args.output else None
    lang = _get_lang(args.lang)

    try:
        generate_ppt(md_path, output_path, lang, dry_run=args.dry_run)
    except SystemExit:
        return 1
    except Exception as exc:
        print(f"[agent4ppt] Unexpected error: {exc}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1

    return 0


if __name__ == "__main__":
    sys.exit(main())
